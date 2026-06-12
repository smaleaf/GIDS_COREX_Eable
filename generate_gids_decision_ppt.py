#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
GIDS Corex 适配方案决策 PPT
面向领导汇报、定方案
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ─── 配色系统 ───────────────────────────────────────────────────
BG_WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT_CYAN   = RGBColor(0x51, 0xB9, 0xD6)
TITLE_NAVY    = RGBColor(0x1B, 0x2F, 0x4A)
MID_BLUE      = RGBColor(0x2E, 0x86, 0xC1)
BODY_GRAY     = RGBColor(0x4A, 0x5A, 0x6A)
LIGHT_GRAY    = RGBColor(0x8A, 0x9B, 0xAA)
SECTION_DARK  = RGBColor(0x0D, 0x1F, 0x35)
GREEN         = RGBColor(0x1E, 0x8B, 0x4C)
GREEN_LIGHT   = RGBColor(0xE8, 0xF8, 0xF0)
ORANGE        = RGBColor(0xE6, 0x7E, 0x22)
ORANGE_LIGHT  = RGBColor(0xFE, 0xF5, 0xE7)
RED           = RGBColor(0xC0, 0x39, 0x2B)
RED_LIGHT     = RGBColor(0xFD, 0xED, 0xEC)
YELLOW        = RGBColor(0xF3, 0x9C, 0x12)
YELLOW_LIGHT  = RGBColor(0xFE, 0xF9, 0xE7)
CYAN_LIGHT    = RGBColor(0xE8, 0xF6, 0xFD)
CYAN_DARK     = RGBColor(0x1A, 0x7A, 0x9A)
BORDER        = RGBColor(0xD5, 0xE8, 0xF5)
DARK_LINE     = RGBColor(0x2C, 0x3E, 0x50)

SLIDE_W = 12191365
SLIDE_H = 6858000


# ─── 基础绘图函数 ───────────────────────────────────────────────
def rect(slide, x, y, w, h, fill=None, line=None, lw=Pt(1)):
    s = slide.shapes.add_shape(1, Emu(x), Emu(y), Emu(w), Emu(h))
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line:
        s.line.color.rgb = line; s.line.width = lw
    else:
        s.line.fill.background()
    return s


def txt(slide, x, y, w, h, text, size=Pt(13), bold=False,
        color=BODY_GRAY, align=PP_ALIGN.LEFT, font="微软雅黑", wrap=True):
    tb = slide.shapes.add_textbox(Emu(x), Emu(y), Emu(w), Emu(h))
    tf = tb.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = size; r.font.bold = bold
    r.font.color.rgb = color; r.font.name = font
    return tb


def multiline(slide, x, y, w, h, lines, default_size=Pt(12),
              default_color=BODY_GRAY, font="微软雅黑",
              spacing=None, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Emu(x), Emu(y), Emu(w), Emu(h))
    tf = tb.text_frame; tf.word_wrap = True
    for i, line in enumerate(lines):
        if isinstance(line, tuple):
            text, size, bold, color = line
        else:
            text, size, bold, color = line, default_size, False, default_color
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if spacing:
            from pptx.util import Pt as PT
            p.space_after = PT(spacing)
        r = p.add_run(); r.text = text
        r.font.size = size; r.font.bold = bold
        r.font.color.rgb = color; r.font.name = font


def title_bar(slide, title, subtitle=None, bar_color=ACCENT_CYAN):
    """通用标题栏"""
    rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=BG_WHITE)
    rect(slide, 0, 0, SLIDE_W, 480000, fill=bar_color)
    rect(slide, 0, 6380000, SLIDE_W, 60000, fill=bar_color)
    txt(slide, 380000, 80000, SLIDE_W - 760000, 330000, title,
        size=Pt(22), bold=True, color=BG_WHITE)
    if subtitle:
        txt(slide, 380000, 6400000, SLIDE_W - 760000, 250000, subtitle,
            size=Pt(10), color=RGBColor(0x90, 0xCA, 0xE0), align=PP_ALIGN.RIGHT)


def section_page(slide, num, title, sub="", num_color=ACCENT_CYAN):
    rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=SECTION_DARK)
    rect(slide, 0, 0, 200000, SLIDE_H, fill=num_color)
    rect(slide, 200000, SLIDE_H//2 - 5000, SLIDE_W - 200000, 8000, fill=RGBColor(0x25,0x3A,0x52))
    txt(slide, 400000, 1800000, 4000000, 1500000, num,
        size=Pt(96), bold=True, color=num_color, align=PP_ALIGN.LEFT)
    txt(slide, 400000, 3200000, SLIDE_W - 800000, 900000, title,
        size=Pt(38), bold=True, color=BG_WHITE)
    if sub:
        txt(slide, 400000, 4100000, SLIDE_W - 800000, 500000, sub,
            size=Pt(17), color=RGBColor(0x7A, 0xB3, 0xD0))


def badge(slide, x, y, text, bg_color, text_color=BG_WHITE, w=2000000, h=320000):
    rect(slide, x, y, w, h, fill=bg_color)
    txt(slide, x+60000, y+50000, w-120000, h-100000, text,
        size=Pt(12), bold=True, color=text_color, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════
prs = Presentation()
prs.slide_width = Emu(SLIDE_W)
prs.slide_height = Emu(SLIDE_H)
BLANK = prs.slide_layouts[6]


# ══════════════════════════════════════════════════════
# Slide 1  封面
# ══════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=SECTION_DARK)
rect(s, 0, 0, SLIDE_W, 8000, fill=ACCENT_CYAN)
rect(s, 0, SLIDE_H-8000, SLIDE_W, 8000, fill=ACCENT_CYAN)
rect(s, 0, 0, 220000, SLIDE_H, fill=ACCENT_CYAN)

# 标题区
rect(s, 220000, 1600000, SLIDE_W-220000, 3400000, fill=RGBColor(0x12,0x28,0x40))
rect(s, 220000, 1600000, 12000, 3400000, fill=ACCENT_CYAN)

txt(s, 500000, 1750000, SLIDE_W-900000, 650000,
    "GIDS × Corex 适配方案决策",
    size=Pt(38), bold=True, color=BG_WHITE)
txt(s, 500000, 2450000, SLIDE_W-900000, 500000,
    "GPU-Initiated Direct Storage for GNN Training on Iluvatar Platform",
    size=Pt(18), color=RGBColor(0x7A,0xC8,0xE8))
rect(s, 500000, 2980000, 3000000, 6000, fill=ACCENT_CYAN)
txt(s, 500000, 3100000, SLIDE_W-900000, 400000,
    "汇报目的：确定 bam_ptr 路线 vs cuFile 路线，明确工作计划与资源需求",
    size=Pt(14), color=RGBColor(0xA0,0xC8,0xDC))

# 信息栏
info = [
    ("汇报人", "GNN 存储团队"),
    ("日期", "2026-06-12"),
    ("版本", "决策版 v1.0"),
    ("状态", "待定方案"),
]
for i, (k, v) in enumerate(info):
    x = 500000 + i * 2900000
    rect(s, x, 5200000, 2700000, 700000, fill=RGBColor(0x1A,0x35,0x55),
         line=RGBColor(0x2E,0x5A,0x8A), lw=Pt(1.5))
    txt(s, x+120000, 5270000, 2500000, 260000, k, size=Pt(11),
        color=RGBColor(0x7A,0xB0,0xD0))
    txt(s, x+120000, 5520000, 2500000, 310000, v, size=Pt(14),
        bold=True, color=BG_WHITE)


# ══════════════════════════════════════════════════════
# Slide 2  摘要：核心结论（先放结论）
# ══════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
title_bar(s, "核心结论 — Executive Summary", "先看结论再看分析")

txt(s, 380000, 630000, SLIDE_W-760000, 330000,
    "调研结论：bam_ptr 移植技术上可行，但存在一个关键硬件确认项需在本次会议决策",
    size=Pt(14), bold=False, color=BODY_GRAY)

items = [
    (GREEN,  "✅ 已确认可行",
     "<<<>>> 语法",
     "Corex ixc 编译器原生支持 <<<>>> kernel launch 语法（host_runtime.h 已验证），\n无需改写为 ixLaunchKernel()，之前的顾虑不成立。"),
    (GREEN,  "✅ 已确认可行",
     "simt::atomic 替换",
     "BaM 依赖 NVIDIA CCCL simt::atomic，Corex 有完全等价的 cuda::atomic\n（cuda/std/atomic，已确认 thread_scope_device 支持），全局搜索替换即可。"),
    (ORANGE, "✅ 有路径，需 1~2 周",
     "内核模块 GPU DMA（nv-p2p.h）",
     "BaM 内核模块用 NVIDIA nv-p2p.h 做 GPU 内存 DMA 注册。\nCorex 有 gdrapi.h（p2p_token 机制），适配路径存在，需要实现和测试。"),
    (RED,    "⚠️ 关键未知项  需本次决策",
     "GPU kernel 直写 NVMe SQ MMIO",
     "bam_ptr 的核心路径：GPU warp 直接写入 NVMe 控制器的 BAR 寄存器（PCIe 设备写）。\n这要求 IX GPU 支持 PCIe peer write to NVMe BAR。目前未验证，需硬件团队确认。"),
]

for i, (color, status, title_t, desc) in enumerate(items):
    y = 1100000 + i * 1350000
    rect(s, 380000, y, SLIDE_W-760000, 1230000,
         fill=RGBColor(0xF8,0xFC,0xFF), line=color, lw=Pt(2.5))
    rect(s, 380000, y, 380000, 1230000, fill=color)
    # 状态标
    rect(s, 880000, y+100000, 2000000, 280000, fill=color)
    txt(s, 900000, y+130000, 1960000, 230000, status,
        size=Pt(11), bold=True, color=BG_WHITE, align=PP_ALIGN.CENTER)
    txt(s, 880000, y+430000, 3200000, 320000, title_t,
        size=Pt(15), bold=True, color=color)
    txt(s, 880000, y+770000, SLIDE_W-1500000, 400000, desc,
        size=Pt(12), color=BODY_GRAY)


# ══════════════════════════════════════════════════════
# Slide 3  章节分隔 - 背景
# ══════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
section_page(s, "01", "背景：GIDS 要做什么？",
             "Why GIDS / What needs to run on Corex")


# ══════════════════════════════════════════════════════
# Slide 4  GIDS 是什么 + 为什么要在 Corex 上运行
# ══════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
title_bar(s, "GIDS 是什么 / 为何要在 Corex 上运行")

# 左：是什么
rect(s, 380000, 620000, 5400000, 5700000, fill=CYAN_LIGHT, line=ACCENT_CYAN, lw=Pt(2))
txt(s, 530000, 700000, 5200000, 350000, "GIDS 是什么",
    size=Pt(16), bold=True, color=CYAN_DARK)
multiline(s, 530000, 1100000, 5100000, 4900000, [
    ("GPU-Initiated Direct Storage Accesses", Pt(14), True, TITLE_NAVY),
    ("专为超大规模 GNN 训练设计的数据加载系统", Pt(12), False, BODY_GRAY),
    ("", Pt(11), False, BODY_GRAY),
    ("核心价值：", Pt(13), True, CYAN_DARK),
    ("• GPU 直接从 NVMe SSD 读取节点特征数据", Pt(12), False, BODY_GRAY),
    ("• 彻底绕过 CPU，消除数据流水线瓶颈", Pt(12), False, BODY_GRAY),
    ("• 支持 2TB+ 级别图特征（IGB-Full 547M 节点）", Pt(12), False, BODY_GRAY),
    ("", Pt(11), False, BODY_GRAY),
    ("技术原理：", Pt(13), True, CYAN_DARK),
    ("• BaM 框架：GPU kernel 直接写 NVMe 命令队列", Pt(12), False, BODY_GRAY),
    ("• GPU 端 page cache（显存作热数据二级缓存）", Pt(12), False, BODY_GRAY),
    ("• Window Buffer 预取 + Accumulator 合并 I/O", Pt(12), False, BODY_GRAY),
    ("", Pt(11), False, BODY_GRAY),
    ("性能收益（论文数据）：", Pt(13), True, CYAN_DARK),
    ("• 数据加载吞吐 vs mmap：+200~500%", Pt(12), False, GREEN),
    ("• 端到端训练时间：缩短 40~60%", Pt(12), False, GREEN),
])

# 右：为什么 Corex
rect(s, 6100000, 620000, 5700000, 5700000, fill=ORANGE_LIGHT, line=ORANGE, lw=Pt(2))
txt(s, 6250000, 700000, 5500000, 350000, "为什么要适配 Corex 平台",
    size=Pt(16), bold=True, color=ORANGE)
multiline(s, 6250000, 1100000, 5400000, 4900000, [
    ("业务需求：", Pt(13), True, ORANGE),
    ("• Iluvatar GPU 是公司主推算力平台", Pt(12), False, BODY_GRAY),
    ("• GNN 训练是核心应用场景之一", Pt(12), False, BODY_GRAY),
    ("• 需要 GIDS 的加速能力跑在 Corex 上", Pt(12), False, BODY_GRAY),
    ("", Pt(11), False, BODY_GRAY),
    ("技术现状：", Pt(13), True, ORANGE),
    ("• 原始 GIDS 依赖 NVIDIA CUDA + BaM", Pt(12), False, BODY_GRAY),
    ("• Corex 提供 CUDA 兼容层（CUDA API 全覆盖）", Pt(12), False, BODY_GRAY),
    ("• 但 BaM 的 NVMe 底层机制需要确认", Pt(12), False, BODY_GRAY),
    ("", Pt(11), False, BODY_GRAY),
    ("目标：", Pt(13), True, ORANGE),
    ("• GIDS 完整功能在 Corex GPU + NVMe 上运行", Pt(12), False, BODY_GRAY),
    ("• 保留原论文的性能优势", Pt(12), False, BODY_GRAY),
    ("• 支持 IGB、OGB 等大规模图数据集训练", Pt(12), False, BODY_GRAY),
    ("", Pt(11), False, BODY_GRAY),
    ("当前里程碑：", Pt(13), True, ORANGE),
    ("• ✅ cuFile 替代方案已移植完成（保底路线）", Pt(12), False, GREEN),
    ("• ⚠️ bam_ptr 原生路线可行性待确认", Pt(12), False, ORANGE),
])


# ══════════════════════════════════════════════════════
# Slide 5  章节分隔 - 两条路线
# ══════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
section_page(s, "02", "两条技术路线对比",
             "cuFile Route vs bam_ptr Route — Trade-offs")


# ══════════════════════════════════════════════════════
# Slide 6  两条路线全景对比
# ══════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
title_bar(s, "两条路线全景对比", "cuFile（保底已有）vs bam_ptr（原生移植）")

# 中间决策轴
rect(s, SLIDE_W//2 - 30000, 600000, 60000, 5700000,
     fill=RGBColor(0xCC,0xCC,0xCC))
txt(s, SLIDE_W//2 - 400000, 580000, 800000, 280000, "VS",
    size=Pt(22), bold=True, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# ── 左侧：路线A cuFile ──
rect(s, 250000, 620000, 5500000, 5750000, fill=GREEN_LIGHT, line=GREEN, lw=Pt(2.5))
rect(s, 250000, 620000, 5500000, 430000, fill=GREEN)
txt(s, 400000, 650000, 5200000, 380000,
    "路线 A  ✅  cuFile + IXFeatureStore（已完成）",
    size=Pt(15), bold=True, color=BG_WHITE)

route_a = [
    ("架构变化：", Pt(13), True, GREEN),
    ("  BAM_Feature_Store → IXFeatureStore", Pt(12), False, BODY_GRAY),
    ("  bam_ptr.read() → cuFileRead() GPU DMA", Pt(12), False, BODY_GRAY),
    ("  /dev/libnvmX → /mnt/nvme0/node_feat.bin", Pt(12), False, BODY_GRAY),
    ("  GPU page cache → OS kernel page cache", Pt(12), False, BODY_GRAY),
    ("", Pt(11), False, BODY_GRAY),
    ("优势：", Pt(13), True, GREEN),
    ("  ✅ 代码已写完，sandbox + Docker 编译通过", Pt(12), False, GREEN),
    ("  ✅ 依赖 Corex 官方 GDS 栈（itrfs.ko 已验证）", Pt(12), False, GREEN),
    ("  ✅ 无硬件未知风险，可快速端到端验证", Pt(12), False, GREEN),
    ("  ✅ 自动 POSIX fallback，程序不会崩溃", Pt(12), False, GREEN),
    ("", Pt(11), False, BODY_GRAY),
    ("代价：", Pt(13), True, ORANGE),
    ("  ⚠ GPU kernel 内无透明 page fault 语义", Pt(12), False, ORANGE),
    ("  ⚠ CPU 发起 DMA，多一次调度延迟", Pt(12), False, ORANGE),
    ("  ⚠ 理论性能低于 bam_ptr（预估 20~40%）", Pt(12), False, ORANGE),
    ("", Pt(11), False, BODY_GRAY),
    ("当前状态：", Pt(13), True, GREEN),
    ("  核心代码完成，DGL 编译进行中", Pt(12), False, BODY_GRAY),
    ("  预计 1~2 周内可端到端跑通", Pt(12), False, BODY_GRAY),
]
multiline(s, 380000, 1130000, 5250000, 5050000, route_a)

# ── 右侧：路线B bam_ptr ──
rect(s, 6250000, 620000, 5700000, 5750000, fill=CYAN_LIGHT, line=MID_BLUE, lw=Pt(2.5))
rect(s, 6250000, 620000, 5700000, 430000, fill=MID_BLUE)
txt(s, 6400000, 650000, 5400000, 380000,
    "路线 B  ⚠️  bam_ptr 原生适配（待评估）",
    size=Pt(15), bold=True, color=BG_WHITE)

route_b = [
    ("架构变化：", Pt(13), True, MID_BLUE),
    ("  保留 bam_ptr<T> + page_cache_t 语义", Pt(12), False, BODY_GRAY),
    ("  simt::atomic → cuda::atomic（简单替换）", Pt(12), False, BODY_GRAY),
    ("  nv-p2p.h → Corex gdrapi.h（内核模块改造）", Pt(12), False, BODY_GRAY),
    ("  GPU 直写 NVMe BAR（需硬件验证）", Pt(12), False, BODY_GRAY),
    ("", Pt(11), False, BODY_GRAY),
    ("优势：", Pt(13), True, MID_BLUE),
    ("  ✅ 完全保留原论文 GPU-Direct 路径", Pt(12), False, MID_BLUE),
    ("  ✅ 性能上限与 NVIDIA 版 GIDS 一致", Pt(12), False, MID_BLUE),
    ("  ✅ 无 CPU 调度开销，延迟最低", Pt(12), False, MID_BLUE),
    ("  ✅ bam_ptr API 对上层完全透明", Pt(12), False, MID_BLUE),
    ("", Pt(11), False, BODY_GRAY),
    ("风险：", Pt(13), True, RED),
    ("  ❌ IX GPU PCIe peer write 到 NVMe BAR 未验证", Pt(12), False, RED),
    ("  ❌ 内核模块需改造（涉及驱动团队）", Pt(12), False, ORANGE),
    ("  ❌ 工期不确定（2~6 周，视硬件支持结论）", Pt(12), False, ORANGE),
    ("", Pt(11), False, BODY_GRAY),
    ("当前状态：", Pt(13), True, MID_BLUE),
    ("  simt::atomic 替换方案已验证可行", Pt(12), False, BODY_GRAY),
    ("  内核 nv-p2p → gdrapi 路径待实现", Pt(12), False, BODY_GRAY),
]
multiline(s, 6380000, 1130000, 5450000, 5050000, route_b)


# ══════════════════════════════════════════════════════
# Slide 7  章节分隔 - bam_ptr 适配详解
# ══════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
section_page(s, "03", "bam_ptr 适配的三层技术障碍",
             "Layer-by-layer technical analysis for bam_ptr porting")


# ══════════════════════════════════════════════════════
# Slide 8  三层障碍详解（横向布局）
# ══════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
title_bar(s, "bam_ptr 适配三层障碍逐层分析",
          "Green = solvable, Red = requires HW team confirmation")

layers = [
    ("层 1", "simt::atomic", GREEN, "★☆☆  1~2 天",
     [
         "BaM 所有 page_cache_t / nvm_parallel_queue",
         "大量使用 NVIDIA CCCL simt::atomic",
         "simt::thread_scope_device / memory_order",
     ],
     [
         "Corex 有 cuda/std/atomic（已确认！）",
         "thread_scope_device 完全支持",
         "fetch_add/fetch_sub/store/load 一一对应",
         "一条 sed 全局替换即可",
     ],
     "✅ 立即可做",
     """// 修复命令（30秒完成）
sed -i 's/simt::/cuda::/g' bam/include/*.h
sed -i 's/<simt\\/atomic>/<cuda\\/std\\/atomic>/g' \\
    bam/include/*.h"""),

    ("层 2", "libnvm.ko GPU DMA\nnv-p2p.h → gdrapi.h", ORANGE, "★★☆  1~2 周",
     [
         "BaM 内核模块用 nvidia_p2p_get_pages()",
         "将 GPU 物理内存页注册为 NVMe DMA 目标",
         "nv-p2p.h 是 NVIDIA 专有内核 API",
     ],
     [
         "Corex 有 gdrapi.h（GDR 用户态 API）",
         "gdr_pin_buffer() + p2p_token 机制",
         "IX_POINTER_ATTRIBUTE_P2P_TOKENS 可获 token",
         "内核模块适配或改为用户态路径",
     ],
     "⚠️ 有路径需实现",
     """// Corex 等价路径
ixdrvPointerGetAttributes(
  IXDRV_POINTER_ATTRIBUTE_P2P_TOKENS,
  &tokens, gpu_ptr);
gdr_pin_buffer(g, addr, size,
  tokens.p2pToken,
  tokens.va_space, &handle);"""),

    ("层 3", "GPU kernel 直写\nNVMe SQ MMIO", RED, "★★★  未知",
     [
         "bam_ptr cache miss 时 GPU warp 写入",
         "NVMe SQ 寄存器（PCIe BAR 空间）",
         "要求 NVMe BAR 映射进 GPU VA 空间",
         "即 PCIe device-to-device peer write",
     ],
     [
         "NVIDIA: GPUDirect RDMA 支持此路径",
         "Corex IX GPU: 尚未验证！",
         "需向 Iluvatar HW/驱动团队确认",
         "若支持：bam_ptr 完整适配可行",
         "若不支持：路线 A（cuFile）为唯一可行方案",
     ],
     "❌ 本次决策关键点",
     """// 需要确认的问题
ixDeviceCanAccessPeer() 覆盖 GPU-GPU
但 GPU → NVMe BAR 属于不同 PCIe 设备
目前无文档/API 确认 IX GPU 支持
需硬件团队给出明确答复"""),
]

col_w = (SLIDE_W - 800000) // 3
for i, (num, name, color, effort, prob, sol, verdict, code) in enumerate(layers):
    x = 350000 + i * (col_w + 50000)
    y = 620000

    # 卡片
    rect(s, x, y, col_w, 5700000,
         fill=RGBColor(0xF8,0xFC,0xFF), line=color, lw=Pt(2.5))
    rect(s, x, y, col_w, 400000, fill=color)

    # 顶部
    txt(s, x+80000, y+50000, 800000, 310000, num,
        size=Pt(22), bold=True, color=BG_WHITE)
    txt(s, x+900000, y+60000, col_w-1000000, 180000,
        effort, size=Pt(12), bold=True, color=BG_WHITE)
    txt(s, x+80000, y+250000, col_w-160000, 160000,
        name, size=Pt(13), bold=True, color=BG_WHITE, wrap=False)

    # 问题
    txt(s, x+80000, y+470000, col_w-160000, 220000,
        "问题：", size=Pt(12), bold=True, color=RED)
    prob_lines = [("• "+p, Pt(11), False, BODY_GRAY) for p in prob]
    multiline(s, x+80000, y+690000, col_w-160000, 1200000, prob_lines)

    # 解决方案
    txt(s, x+80000, y+1980000, col_w-160000, 220000,
        "解法：", size=Pt(12), bold=True, color=GREEN)
    sol_lines = [("• "+p, Pt(11), False, BODY_GRAY) for p in sol]
    multiline(s, x+80000, y+2200000, col_w-160000, 1400000, sol_lines)

    # 结论徽章
    rect(s, x+80000, y+3750000, col_w-160000, 300000, fill=color)
    txt(s, x+80000, y+3790000, col_w-160000, 230000, verdict,
        size=Pt(12), bold=True, color=BG_WHITE, align=PP_ALIGN.CENTER)

    # 代码
    rect(s, x+80000, y+4120000, col_w-160000, 1500000,
         fill=RGBColor(0xF0,0xF0,0xF0), line=RGBColor(0xCC,0xCC,0xCC))
    multiline(s, x+130000, y+4170000, col_w-260000, 1380000,
              [(code, Pt(10), False, RGBColor(0x2E,0x60,0x9A))],
              font="Consolas")


# ══════════════════════════════════════════════════════
# Slide 9  章节分隔 - 关键验证
# ══════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
section_page(s, "04", "关键验证项 — 决策前必须回答",
             "The one question that determines everything", num_color=RED)


# ══════════════════════════════════════════════════════
# Slide 10  关键问题：IX GPU 能否写 NVMe BAR
# ══════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
title_bar(s, "唯一关键验证项", "答案决定走哪条路线", bar_color=RED)

# 问题框
rect(s, 380000, 600000, SLIDE_W-760000, 1100000,
     fill=RED_LIGHT, line=RED, lw=Pt(3))
rect(s, 380000, 600000, 20000, 1100000, fill=RED)
txt(s, 600000, 680000, SLIDE_W-1200000, 320000,
    "核心问题（需向 Iluvatar 硬件/驱动团队确认）",
    size=Pt(14), bold=True, color=RED)
txt(s, 600000, 1020000, SLIDE_W-1200000, 500000,
    "IX GPU kernel 中，能否通过 PCIe 直接写入 NVMe 控制器的 BAR 寄存器（MMIO）？\n即：NVMe SQ 寄存器地址能否映射到 GPU 虚拟地址空间，并允许 GPU warp store 写入？",
    size=Pt(16), bold=True, color=TITLE_NAVY)

# 两种答案的后续
txt(s, 380000, 1900000, SLIDE_W-760000, 280000,
    "两种答案对应完全不同的工作计划：",
    size=Pt(13), bold=True, color=BODY_GRAY)

# 答案 YES
rect(s, 380000, 2250000, 5500000, 3900000,
     fill=GREEN_LIGHT, line=GREEN, lw=Pt(2.5))
rect(s, 380000, 2250000, 5500000, 420000, fill=GREEN)
txt(s, 530000, 2280000, 5200000, 370000,
    "答案：YES（IX GPU 支持此 PCIe peer write）",
    size=Pt(15), bold=True, color=BG_WHITE)
yes_lines = [
    ("采用路线 B：bam_ptr 完整移植", Pt(14), True, GREEN),
    ("", Pt(11), False, BODY_GRAY),
    ("工作步骤：", Pt(12), True, TITLE_NAVY),
    ("  Week 1：simt::atomic → cuda::atomic（已有方案）", Pt(12), False, BODY_GRAY),
    ("  Week 1-2：编译 libnvm.ko 内核模块（替换 nv-p2p.h）", Pt(12), False, BODY_GRAY),
    ("  Week 2-3：验证 GPU kernel 写 NVMe SQ 路径", Pt(12), False, BODY_GRAY),
    ("  Week 3-4：集成 page_cache_t，端到端跑通", Pt(12), False, BODY_GRAY),
    ("  Week 5-6：性能调优，与 NVIDIA 版对比", Pt(12), False, BODY_GRAY),
    ("", Pt(11), False, BODY_GRAY),
    ("预期工期：4~6 周", Pt(13), True, GREEN),
    ("性能预期：与 NVIDIA 原版 GIDS 相当", Pt(12), False, GREEN),
]
multiline(s, 530000, 2750000, 5200000, 3200000, yes_lines)

# 答案 NO
rect(s, 6300000, 2250000, 5500000, 3900000,
     fill=ORANGE_LIGHT, line=ORANGE, lw=Pt(2.5))
rect(s, 6300000, 2250000, 5500000, 420000, fill=ORANGE)
txt(s, 6450000, 2280000, 5200000, 370000,
    "答案：NO（IX GPU 不支持此 PCIe peer write）",
    size=Pt(15), bold=True, color=BG_WHITE)
no_lines = [
    ("采用路线 A：cuFile + IXFeatureStore（已有）", Pt(14), True, ORANGE),
    ("", Pt(11), False, BODY_GRAY),
    ("当前状态：已完成 80%，可立即推进", Pt(12), False, BODY_GRAY),
    ("", Pt(11), False, BODY_GRAY),
    ("剩余工作：", Pt(12), True, TITLE_NAVY),
    ("  Week 1：完成 DGL CUDA 版编译（进行中）", Pt(12), False, BODY_GRAY),
    ("  Week 1-2：端到端训练跑通验证", Pt(12), False, BODY_GRAY),
    ("  Week 2-4：多 SSD / Window Buffer 优化", Pt(12), False, BODY_GRAY),
    ("  Week 4-6：异构图 + 性能 Profiling", Pt(12), False, BODY_GRAY),
    ("", Pt(11), False, BODY_GRAY),
    ("性能预期：略低于 NVIDIA 原版（20~40%）", Pt(12), False, ORANGE),
    ("预期工期：2~3 周端到端跑通", Pt(13), True, GREEN),
]
multiline(s, 6450000, 2750000, 5200000, 3200000, no_lines)

# 中间箭头
txt(s, SLIDE_W//2 - 250000, 3700000, 500000, 500000, "?",
    size=Pt(60), bold=True, color=RED, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════
# Slide 11  章节分隔 - 工作计划
# ══════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
section_page(s, "05", "推荐工作计划",
             "Recommended action plan regardless of hardware answer")


# ══════════════════════════════════════════════════════
# Slide 12  双轨并行工作计划
# ══════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
title_bar(s, "推荐：双轨并行 — 不赌答案，同步推进")

txt(s, 380000, 600000, SLIDE_W-760000, 300000,
    "建议：在等待硬件团队确认的同时，两条线同步推进，确保无论答案如何都不浪费工时",
    size=Pt(13), color=BODY_GRAY)

# 时间轴
weeks = ["第1周", "第2周", "第3周", "第4周", "第5周", "第6周"]
week_w = (SLIDE_W - 1200000) // 6
header_y = 1100000

rect(s, 600000, header_y, 500000, 380000, fill=TITLE_NAVY)
txt(s, 600000, header_y+80000, 500000, 260000, "任务",
    size=Pt(12), bold=True, color=BG_WHITE, align=PP_ALIGN.CENTER)
for wi, wk in enumerate(weeks):
    wx = 1100000 + wi * (week_w + 10000)
    rect(s, wx, header_y, week_w, 380000, fill=ACCENT_CYAN)
    txt(s, wx+20000, header_y+80000, week_w-40000, 260000, wk,
        size=Pt(12), bold=True, color=BG_WHITE, align=PP_ALIGN.CENTER)

# 任务行
tasks = [
    # (任务名, 颜色, 开始week(0-based), 跨越weeks数, 说明)
    ("硬件确认：IX GPU PCIe peer write", RED, 0, 1, "❶ 本周内提问硬件团队"),
    ("simt::atomic → cuda::atomic 修复", GREEN, 0, 1, "无论哪条路线都需要"),
    ("DGL CUDA 版编译完成（进行中）", GREEN, 0, 1, "两条路线共用"),
    ("路线A: IXFeatureStore 端到端验证", GREEN, 1, 2, "cuFile 保底路线"),
    ("路线A: 性能调优 + 多SSD", GREEN, 3, 2, "Window Buffer / Accumulator"),
    ("路线B: libnvm.ko 内核模块移植", MID_BLUE, 1, 2, "仅 YES 情况下"),
    ("路线B: bam_ptr 端到端集成验证", MID_BLUE, 3, 2, "仅 YES 情况下"),
    ("性能基准测试 + 汇报", TITLE_NAVY, 5, 1, "两条路线均需"),
]

row_h = 530000
for ri, (name, color, start, span, note) in enumerate(tasks):
    ry = header_y + 420000 + ri * (row_h + 15000)
    # 标签
    task_bg = RGBColor(0xF0,0xF8,0xFF) if ri % 2 == 0 else BG_WHITE
    rect(s, 600000, ry, 500000, row_h, fill=task_bg, line=BORDER)
    # 任务条
    bar_x = 1100000 + start * (week_w + 10000)
    bar_w = span * (week_w + 10000) - 10000
    rect(s, bar_x, ry + 60000, bar_w, row_h - 120000, fill=color)
    txt(s, bar_x + 60000, ry + 100000, bar_w - 120000, row_h - 220000,
        name, size=Pt(11), bold=True, color=BG_WHITE)
    # 注释
    note_x = 1100000 + 6 * (week_w + 10000) + 30000
    txt(s, note_x, ry + 130000, 1000000, row_h - 260000, note,
        size=Pt(10), color=color)
    # 行号
    txt(s, 620000, ry + 160000, 460000, row_h - 320000, str(ri+1),
        size=Pt(12), bold=True, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# 图例
legend_y = header_y + 420000 + len(tasks) * (row_h + 15000) + 100000
for c, label in [(GREEN, "两条路线共用 / 路线A"), (MID_BLUE, "仅路线B（条件：YES）"), (RED, "决策节点"), (TITLE_NAVY, "里程碑")]:
    rect(s, 600000, legend_y, 300000, 200000, fill=c)
    txt(s, 950000, legend_y + 30000, 2200000, 160000, label, size=Pt(11), color=BODY_GRAY)
    legend_y = 0  # 只显示第一个（简化）

multiline(s, 600000, 6350000, SLIDE_W-760000, 280000, [
    ("图例：  ", Pt(11), True, BODY_GRAY),
])

# 实际图例
legend_items = [(GREEN,"共用/路线A"), (MID_BLUE,"仅路线B"), (RED,"决策节点"), (TITLE_NAVY,"里程碑")]
for i, (c, lb) in enumerate(legend_items):
    lx = 600000 + i * 3000000
    rect(s, lx, 6400000, 200000, 200000, fill=c)
    txt(s, lx + 260000, 6405000, 2500000, 200000, lb, size=Pt(11), color=BODY_GRAY)


# ══════════════════════════════════════════════════════
# Slide 13  当前已完成工作清单
# ══════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
title_bar(s, "当前已完成工作（不论选哪条路线都有用）",
          "Work already done — valuable for both routes")

done_items = [
    ("✅", GREEN, "GIDS 全源码架构分析完成",
     "gids_kernel.cu / gids_nvme.cu / GIDS.py 深度分析，所有依赖清单梳理完毕"),
    ("✅", GREEN, "bam_ptr 适配可行性分析完成",
     "三层障碍（simt::atomic / nv-p2p.h / NVMe BAR）全部识别，每层均有具体适配方案"),
    ("✅", GREEN, "<<<>>> 语法支持确认",
     "Corex ixc host_runtime.h 验证，原生支持，无需任何修改，消除原顾虑"),
    ("✅", GREEN, "simt::atomic → cuda::atomic 方案确认",
     "Corex cuda/std/atomic 验证，thread_scope_device 完全支持，已有替换命令"),
    ("✅", GREEN, "路线A: IXFeatureStore 核心代码完成",
     "ix_feature_store.cu (~450行) + GIDS_IX.py + build_ix.sh + run.sh 全部就绪"),
    ("✅", GREEN, "基础设施全部就绪",
     "itrfs.ko 已加载 / /dev/itrfs 就绪 / cooperative_groups 补丁部署 / Corex 4.5.0 SDK 全部安装"),
    ("✅", GREEN, "DGL v1.1.3 编译 6 项问题修复",
     "-Xcompiler格式 / fp16.cuh重定义 / CCCL禁用 / omp.h路径 / array_iterator / gpu_cache禁用"),
    ("🔧", MID_BLUE, "DGL CUDA 版编译进行中",
     "主库 ~90% 完成，gpu_cache 禁用绕过，预计本周内完成"),
]

for i, (icon, color, title_t, desc) in enumerate(done_items):
    col = i % 2
    row = i // 2
    x = 380000 + col * 5900000
    y = 700000 + row * 1350000
    rect(s, x, y, 5700000, 1250000,
         fill=GREEN_LIGHT if color == GREEN else CYAN_LIGHT,
         line=color, lw=Pt(1.5))
    txt(s, x + 80000, y + 80000, 500000, 500000, icon,
        size=Pt(22), bold=True, color=color, align=PP_ALIGN.CENTER)
    txt(s, x + 620000, y + 100000, 4930000, 360000, title_t,
        size=Pt(14), bold=True, color=color)
    txt(s, x + 620000, y + 500000, 4930000, 650000, desc,
        size=Pt(11), color=BODY_GRAY)


# ══════════════════════════════════════════════════════
# Slide 14  章节分隔 - 结论与建议
# ══════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
section_page(s, "06", "结论与建议",
             "Recommendation & Decision Points for Leadership")


# ══════════════════════════════════════════════════════
# Slide 15  决策页
# ══════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
title_bar(s, "本次会议需要决策的三件事", "Three decisions needed today")

decisions = [
    ("决策 1", RED, "硬件团队确认任务",
     "由谁负责、何时向 Iluvatar 硬件/驱动团队提出以下问题并回答：",
     ["\"IX GPU kernel 线程能否通过 PCIe 向 NVMe 控制器 BAR 写入数据？\"",
      "\"libnvm 内核模块是否可在 Linux 5.4 + Corex 驱动上编译？\"",
      "\"Corex 是否有内核侧 GPU P2P DMA API（对标 nv-p2p.h）？\""],
     "建议：本周内（2 个工作日）给出答复"),
    ("决策 2", ORANGE, "资源分配",
     "双轨并行需要明确人员分工：",
     ["路线A（cuFile）：由谁负责 DGL 编译 + 端到端验证",
      "路线B（bam_ptr）：由谁负责内核模块移植和 GPU 硬件验证",
      "两条线是否同一人，还是分两组并行"],
     "建议：路线A 1人主导，路线B 需 1 名内核/驱动背景工程师"),
    ("决策 3", MID_BLUE, "里程碑与汇报节点",
     "确定关键时间节点：",
     ["T+3天：硬件团队给出 IX GPU PCIe peer write 答复",
      "T+1周：路线A 端到端训练跑通（DGL 编译完成后）",
      "T+3周：性能基准测试数据（对比 mmap 基线）",
      "T+6周：最终方案确认 + 完整性能报告"],
     "建议：3 天内先拿到硬件答复，再确定总体工期"),
]

for i, (num, color, title_t, sub, points, rec) in enumerate(decisions):
    y = 680000 + i * 1950000
    rect(s, 380000, y, SLIDE_W-760000, 1830000,
         fill=RGBColor(0xF8,0xFB,0xFF), line=color, lw=Pt(2.5))
    rect(s, 380000, y, 600000, 1830000, fill=color)
    txt(s, 440000, y+200000, 490000, 500000, num,
        size=Pt(20), bold=True, color=BG_WHITE, align=PP_ALIGN.CENTER)
    txt(s, 1100000, y+80000, 4000000, 380000, title_t,
        size=Pt(16), bold=True, color=color)
    txt(s, 1100000, y+470000, 4000000, 280000, sub,
        size=Pt(12), color=BODY_GRAY)
    pts_lines = [("  • "+p, Pt(11), False, BODY_GRAY) for p in points]
    multiline(s, 1100000, y+760000, 6500000, 950000, pts_lines)
    # 建议标签
    rect(s, 7800000, y+80000, SLIDE_W-8200000, 350000, fill=color)
    txt(s, 7850000, y+120000, SLIDE_W-8300000, 270000, "建议",
        size=Pt(11), bold=True, color=BG_WHITE, align=PP_ALIGN.CENTER)
    txt(s, 7800000, y+500000, SLIDE_W-8200000, 1250000, rec,
        size=Pt(11), color=color)


# ══════════════════════════════════════════════════════
# Slide 16  封底
# ══════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=SECTION_DARK)
rect(s, 0, 0, SLIDE_W, 8000, fill=ACCENT_CYAN)
rect(s, 0, SLIDE_H-8000, SLIDE_W, 8000, fill=ACCENT_CYAN)
rect(s, 0, 0, 220000, SLIDE_H, fill=ACCENT_CYAN)

txt(s, 500000, 1500000, SLIDE_W-1000000, 800000,
    "一句话总结", size=Pt(20), color=RGBColor(0x7A,0xB8,0xD5))
txt(s, 500000, 2300000, SLIDE_W-1000000, 1200000,
    "路线A（cuFile）已就绪，1~2周可端到端跑通\n路线B（bam_ptr）技术可行，但需先确认 IX GPU 是否支持 PCIe peer write 到 NVMe BAR",
    size=Pt(24), bold=True, color=BG_WHITE)

rect(s, 500000, 3800000, SLIDE_W-1000000, 6000,
     fill=RGBColor(0x2E,0x5A,0x8A))

txt(s, 500000, 4000000, SLIDE_W-1000000, 500000,
    "最优策略：本周内拿到硬件确认，同步推进路线A保底 + 路线B评估",
    size=Pt(16), color=ACCENT_CYAN)

action_items = [
    ("立即行动①", "向 Iluvatar 提问：IX GPU PCIe peer write to NVMe BAR？", RED),
    ("立即行动②", "继续推进 DGL 编译 + 路线A 端到端验证（本周内）", GREEN),
    ("立即行动③", "分配路线B内核模块工程师（等待硬件答复后启动）", MID_BLUE),
]
for i, (label, action, color) in enumerate(action_items):
    rect(s, 500000, 4700000 + i*520000, SLIDE_W-1000000, 480000,
         fill=RGBColor(0x1A,0x35,0x55), line=color, lw=Pt(1.5))
    txt(s, 600000, 4750000 + i*520000, 2000000, 380000, label,
        size=Pt(12), bold=True, color=color)
    txt(s, 2700000, 4750000 + i*520000, SLIDE_W-3300000, 380000, action,
        size=Pt(12), color=BG_WHITE)

txt(s, 500000, 6500000, SLIDE_W-1000000, 280000,
    "GIDS × Corex  |  Iluvatar GPU Platform  |  2026-06-12",
    size=Pt(11), color=RGBColor(0x4A,0x7A,0xA0), align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════
output = "/root/GIDS_cufile/GIDS-Corex适配方案决策-汇报PPT.pptx"
prs.save(output)
print(f"✅ PPT 已生成：{output}")
print(f"   共 {len(prs.slides)} 张幻灯片")
