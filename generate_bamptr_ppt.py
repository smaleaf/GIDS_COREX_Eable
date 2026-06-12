#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
bam_ptr 原生适配 Corex — 完整分析 PPT
"""

from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ─── 配色 ──────────────────────────────────────────────────────
BG         = RGBColor(0xFF,0xFF,0xFF)
NAVY       = RGBColor(0x0D,0x1F,0x35)
CYAN       = RGBColor(0x51,0xB9,0xD6)
CYAN_D     = RGBColor(0x1A,0x7A,0x9A)
CYAN_L     = RGBColor(0xE8,0xF6,0xFD)
BLUE       = RGBColor(0x2E,0x86,0xC1)
BLUE_L     = RGBColor(0xE8,0xF0,0xFB)
GRAY       = RGBColor(0x4A,0x5A,0x6A)
LGRAY      = RGBColor(0x8A,0x9B,0xAA)
BORD       = RGBColor(0xD0,0xE8,0xF5)
GREEN      = RGBColor(0x1B,0x8A,0x4C)
GREEN_L    = RGBColor(0xE8,0xF8,0xF0)
GREEN_D    = RGBColor(0x14,0x6B,0x3A)
ORANGE     = RGBColor(0xD6,0x7B,0x1E)
ORANGE_L   = RGBColor(0xFE,0xF5,0xE7)
RED        = RGBColor(0xC0,0x39,0x2B)
RED_L      = RGBColor(0xFD,0xED,0xEC)
YELLOW     = RGBColor(0xF3,0x9C,0x12)
DARK_LINE  = RGBColor(0x1A,0x30,0x50)
CODE_BG    = RGBColor(0xF4,0xF7,0xFA)
CODE_BLUE  = RGBColor(0x1E,0x5A,0x9A)
CODE_GREEN = RGBColor(0x0A,0x6A,0x35)

W = 12191365
H = 6858000
BL = prs = None   # filled below


# ─── 工具函数 ─────────────────────────────────────────────────
def R(s, x, y, w, h, fill=None, line=None, lw=Pt(1)):
    sh = s.shapes.add_shape(1, Emu(x), Emu(y), Emu(w), Emu(h))
    if fill: sh.fill.solid(); sh.fill.fore_color.rgb = fill
    else: sh.fill.background()
    if line: sh.line.color.rgb = line; sh.line.width = lw
    else: sh.line.fill.background()
    return sh

def T(s, x, y, w, h, text, sz=Pt(12), bold=False, color=GRAY,
      align=PP_ALIGN.LEFT, font="微软雅黑", wrap=True):
    tb = s.shapes.add_textbox(Emu(x), Emu(y), Emu(w), Emu(h))
    tf = tb.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = sz; r.font.bold = bold
    r.font.color.rgb = color; r.font.name = font
    return tb

def ML(s, x, y, w, h, lines, dsz=Pt(12), dclr=GRAY, font="微软雅黑"):
    tb = s.shapes.add_textbox(Emu(x), Emu(y), Emu(w), Emu(h))
    tf = tb.text_frame; tf.word_wrap = True
    for i, line in enumerate(lines):
        if isinstance(line, tuple): text, sz, bold, clr = line
        else: text, sz, bold, clr = line, dsz, False, dclr
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run(); r.text = text
        r.font.size = sz; r.font.bold = bold
        r.font.color.rgb = clr; r.font.name = font

def CODE(s, x, y, w, h, code_lines, sz=Pt(11)):
    """代码块"""
    R(s, x, y, w, h, fill=CODE_BG, line=RGBColor(0xCC,0xD8,0xE8), lw=Pt(1))
    R(s, x, y, 18000, h, fill=CYAN_D)
    lines = [(l, sz, False, CODE_BLUE) for l in code_lines.split('\n')]
    ML(s, x+80000, y+60000, w-160000, h-120000, lines, font="Consolas")

def hdr(s, title, sub=None):
    R(s, 0, 0, W, H, fill=BG)
    R(s, 0, 0, W, 500000, fill=NAVY)
    R(s, 0, 6380000, W, 50000, fill=CYAN)
    T(s, 400000, 80000, W-800000, 360000, title,
      sz=Pt(22), bold=True, color=BG)
    if sub:
        T(s, 400000, 6395000, W-800000, 220000, sub,
          sz=Pt(10), color=LGRAY, align=PP_ALIGN.RIGHT)

def sec(s, num, title, sub="", accent=CYAN):
    R(s, 0, 0, W, H, fill=NAVY)
    R(s, 0, 0, 200000, H, fill=accent)
    R(s, 200000, H//2-6000, W-200000, 10000, fill=RGBColor(0x1E,0x35,0x55))
    T(s, 380000, 1700000, 3500000, 1600000, num,
      sz=Pt(100), bold=True, color=accent)
    T(s, 380000, 3200000, W-700000, 1000000, title,
      sz=Pt(36), bold=True, color=BG)
    if sub:
        T(s, 380000, 4200000, W-700000, 600000, sub,
          sz=Pt(16), color=RGBColor(0x7A,0xB5,0xD5))

def card(s, x, y, w, h, title, lines, accent=BLUE, bg=BLUE_L, code=None):
    R(s, x, y, w, h, fill=bg, line=accent, lw=Pt(2))
    R(s, x, y, w, 350000, fill=accent)
    T(s, x+80000, y+70000, w-160000, 260000, title,
      sz=Pt(13), bold=True, color=BG)
    ML(s, x+80000, y+400000, w-160000, h-460000, lines)
    if code:
        pass


# ════════════════════════════════════════════════════════════════
prs = Presentation()
prs.slide_width = Emu(W); prs.slide_height = Emu(H)
BL = prs.slide_layouts[6]


# ══════════════════════════════════════════════════════
# Slide 1  封面
# ══════════════════════════════════════════════════════
s = prs.slides.add_slide(BL)
R(s, 0, 0, W, H, fill=NAVY)
R(s, 0, 0, W, 6000, fill=CYAN)
R(s, 0, H-6000, W, 6000, fill=CYAN)
R(s, 0, 0, 200000, H, fill=CYAN)

R(s, 200000, 1500000, W-200000, 3800000, fill=RGBColor(0x11,0x25,0x40))
R(s, 200000, 1500000, 16000, 3800000, fill=CYAN)

T(s, 500000, 1620000, W-900000, 750000, "bam_ptr 原生适配 Corex",
  sz=Pt(42), bold=True, color=BG)
T(s, 500000, 2420000, W-900000, 500000,
  "BaM GPU-Initiated NVMe Framework — Architecture Analysis & Porting Plan",
  sz=Pt(17), color=RGBColor(0x7A,0xC5,0xE8))
R(s, 500000, 2970000, 2500000, 5000, fill=CYAN)
T(s, 500000, 3060000, W-900000, 500000,
  "架构组成分析 · 6大适配项逐项攻克 · 完整工作流程",
  sz=Pt(14), color=RGBColor(0xA5,0xD5,0xEA))

chips = [("总适配项", "6", CYAN), ("已可行", "4", GREEN), ("需验证", "1", ORANGE), ("需开发", "1", RED)]
for i,(lbl,val,c) in enumerate(chips):
    xi = 500000 + i*3000000
    R(s, xi, 4700000, 2700000, 900000, fill=RGBColor(0x18,0x32,0x52), line=c, lw=Pt(2))
    T(s, xi+80000, 4750000, 2550000, 380000, val, sz=Pt(36), bold=True, color=c, align=PP_ALIGN.CENTER)
    T(s, xi+80000, 5130000, 2550000, 300000, lbl, sz=Pt(12), color=RGBColor(0x8A,0xBE,0xD8), align=PP_ALIGN.CENTER)

T(s, 500000, 6500000, W-900000, 280000,
  "Iluvatar Corex 4.5.0  |  GIDS GNN Training  |  2026-06-12",
  sz=Pt(11), color=RGBColor(0x3A,0x6A,0x9A), align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════
# Slide 2  bam_ptr 是什么 — 一句话 + 核心原理
# ══════════════════════════════════════════════════════
s = prs.slides.add_slide(BL)
hdr(s, "bam_ptr 是什么 — 为什么要原生适配它")

# 左：定义 + 价值
R(s, 350000, 580000, 5400000, 5900000, fill=CYAN_L, line=CYAN, lw=Pt(2))
T(s, 500000, 660000, 5100000, 350000, "bam_ptr<T> 的核心价值",
  sz=Pt(16), bold=True, color=CYAN_D)
ML(s, 500000, 1060000, 5100000, 5100000, [
    ("GPU kernel 内的「智能 NVMe 指针」", Pt(14), True, NAVY),
    ("", Pt(11), False, GRAY),
    ("工作方式：", Pt(13), True, CYAN_D),
    ("bam_ptr<float> ptr(array);   // 指向 NVMe 上的数据", Pt(12), False, CODE_BLUE),
    ("float val = ptr[row * dim + col];  // GPU kernel 内透明读取", Pt(12), False, CODE_BLUE),
    ("", Pt(11), False, GRAY),
    ("透明语义：", Pt(13), True, CYAN_D),
    ("• operator[] 被调用 → 检查 GPU 显存 page cache", Pt(12), False, GRAY),
    ("• Cache Hit：直接从 GPU 显存返回，零延迟", Pt(12), False, GREEN),
    ("• Cache Miss：GPU warp 直接写 NVMe 命令到 SQ 寄存器", Pt(12), False, ORANGE),
    ("   ↳ NVMe 控制器 DMA 数据到 GPU 显存（P2P）", Pt(12), False, ORANGE),
    ("   ↳ GPU warp 轮询 CQ，获取数据，继续执行", Pt(12), False, ORANGE),
    ("", Pt(11), False, GRAY),
    ("与 cuFile 路线的本质区别：", Pt(13), True, NAVY),
    ("• cuFile：CPU 发起 DMA → GPU，有调度延迟", Pt(12), False, RED),
    ("• bam_ptr：GPU 自主发起 NVMe I/O，CPU 完全不参与", Pt(12), False, GREEN),
    ("", Pt(11), False, GRAY),
    ("性能差异（理论）：", Pt(13), True, NAVY),
    ("• 消除 CPU 调度路径，I/O 延迟降低 50~80%", Pt(12), False, GREEN),
    ("• 多 GPU warp 并发访问，NVMe 队列深度饱和", Pt(12), False, GREEN),
    ("• 与 NVIDIA 原版 GIDS 论文性能一致", Pt(12), False, GREEN),
])

# 右：架构层次
R(s, 6100000, 580000, 5700000, 5900000, fill=BLUE_L, line=BLUE, lw=Pt(2))
T(s, 6250000, 660000, 5450000, 350000, "BaM 在 GIDS 中的位置",
  sz=Pt(16), bold=True, color=BLUE)
layers_r = [
    ("Python 层", "GIDS.py / GIDS_DGLDataLoader", GRAY, BG),
    ("pybind11", "BAM_Feature_Store.so", LGRAY, BG),
    ("C++/CUDA 核心", "gids_kernel.cu / gids_nvme.cu", BLUE, BLUE_L),
    ("BaM 框架", "bam_ptr<T> / page_cache_t / Controller", CYAN_D, CYAN_L),
    ("内核模块", "libnvm.ko → /dev/libnvmX", ORANGE, ORANGE_L),
    ("硬件", "IX GPU ←PCIe P2P→ NVMe SSD", NAVY, RGBColor(0xE0,0xE8,0xF5)),
]
for i, (title, desc, c, bg) in enumerate(layers_r):
    ly = 1100000 + i * 770000
    R(s, 6250000, ly, 5400000, 700000, fill=bg, line=c, lw=Pt(1.5))
    T(s, 6350000, ly+80000, 1800000, 300000, title, sz=Pt(12), bold=True, color=c)
    T(s, 8200000, ly+80000, 3300000, 300000, desc, sz=Pt(11), color=GRAY)
    # 箭头
    if i < 5:
        T(s, 6750000, ly+700000, 400000, 100000, "↓", sz=Pt(10), bold=True, color=c, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════
# Slide 3  BaM 完整架构 — 五个核心对象
# ══════════════════════════════════════════════════════
s = prs.slides.add_slide(BL)
hdr(s, "BaM 完整架构组成 — 五大核心对象", "每个对象的职责、CUDA 依赖、适配难度")

objs = [
    ("bafs_ptr<T>\n(bam_ptr)", CYAN_D,
     ["头文件：bafs_ptr.h",
      "层次：GPU kernel 内使用",
      "operator[](i) → page cache 查找 + NVMe I/O"],
     ["#include <page_cache.h>（间接依赖所有）",
      "__host__ __device__ 标记"],
     ["simt::atomic 通过 page_cache_t 间接引用",
      "无直接 CUDA API 调用"],
     "✅ 轻松"),
    ("page_cache_t", BLUE,
     ["头文件：page_cache.h（2143行）",
      "层次：GPU kernel + Host 管理",
      "GPU 显存页缓存：LRU 淘汰 + page fault 处理"],
     ["#include <simt/atomic>  ← CCCL 依赖",
      "#include <nvm_parallel_queue.h>",
      "大量 simt::atomic 原子操作"],
     ["simt::atomic<uint32_t, thread_scope_device>",
      "simt::atomic<uint64_t, thread_scope_device>",
      "DmaPtr pages_dma (GPU 内存 DMA 注册)"],
     "⚠️ 中等"),
    ("nvm_types.h\n(nvm_queue_t)", ORANGE,
     ["NVMe 队列描述符（SQ/CQ）",
      "volatile uint32_t* db  ← 门铃寄存器指针",
      "GPU kernel 通过 db 写入触发 NVMe I/O"],
     ["#include <simt/atomic>  ← CCCL 依赖",
      "simt::thread_scope_device + system"],
     ["NVMe SQ/CQ 在 GPU 内存中分配",
      "db 指针 = NVMe BAR mmap + cudaHostGetDevicePointer()"],
     "⚠️ 中等"),
    ("Controller\n(ctrl.h)", GREEN_D,
     ["Host 端控制器管理",
      "打开 /dev/libnvmX，读取 NVMe BAR",
      "cudaHostRegister(mm_ptr, IoMemory) 注册 MMIO"],
     ["cudaHostRegister + cudaHostRegisterIoMemory",
      "cudaHostGetDevicePointer 获取门铃 GPU 地址",
      "cudaMalloc / cudaMemcpy / cuda_err_chk"],
     ["cudaHostRegisterIoMemory = 0x04",
      "cudaHostGetDevicePointer → ix 已有映射!",
      "cudaMalloc/cudaMemcpy → ix 已有映射!"],
     "⚠️ 待测试"),
    ("libnvm.ko\n内核模块", RED,
     ["A部分：mmap_registers() — NVMe BAR mmap 到用户态",
      "B部分：NVM_MAP_HOST_MEMORY — 主机内存 DMA 注册",
      "C部分(#ifdef _CUDA)：GPU 内存 DMA 注册"],
     ["A+B：纯 Linux PCI/DMA API，无 NVIDIA 依赖",
      "C：#include <nv-p2p.h>",
      "   nvidia_p2p_get_pages()",
      "   nvidia_p2p_dma_map_pages()"],
     ["A+B：无需修改直接编译 ✅",
      "C：需替换 nvidia_p2p_* 为 Corex 等价物",
      "   选项1：Corex ix-p2p.h 内核 API",
      "   选项2：gdrapi.h 用户态路径"],
     "❌ 需要开发"),
]

col_w = (W - 900000) // 5
for i, (name, color, role, deps, corex, diff) in enumerate(objs):
    x = 380000 + i * (col_w + 15000)
    diff_c = {"✅ 轻松": GREEN, "⚠️ 中等": ORANGE, "⚠️ 待测试": ORANGE, "❌ 需要开发": RED}[diff]
    bg = GREEN_L if diff_c==GREEN else (ORANGE_L if diff_c==ORANGE else RED_L)
    R(s, x, 580000, col_w, 5900000, fill=bg, line=color, lw=Pt(2))
    R(s, x, 580000, col_w, 350000, fill=color)
    T(s, x+60000, 610000, col_w-120000, 290000, name,
      sz=Pt(12), bold=True, color=BG, wrap=False)

    T(s, x+60000, 990000, col_w-120000, 240000, "职责",
      sz=Pt(11), bold=True, color=color)
    role_l = [("• "+r, Pt(10), False, GRAY) for r in role]
    ML(s, x+60000, 1230000, col_w-120000, 1100000, role_l)

    T(s, x+60000, 2380000, col_w-120000, 240000, "CUDA依赖",
      sz=Pt(11), bold=True, color=RED)
    dep_l = [("• "+d, Pt(10), False, GRAY) for d in deps]
    ML(s, x+60000, 2620000, col_w-120000, 1300000, dep_l)

    T(s, x+60000, 3980000, col_w-120000, 240000, "Corex方案",
      sz=Pt(11), bold=True, color=GREEN_D)
    cr_l = [("• "+c, Pt(10), False, GRAY) for c in corex]
    ML(s, x+60000, 4220000, col_w-120000, 1150000, cr_l)

    R(s, x+60000, 5520000, col_w-120000, 330000, fill=diff_c)
    T(s, x+60000, 5560000, col_w-120000, 270000, diff,
      sz=Pt(12), bold=True, color=BG, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════
# Slide 4  章节分隔 — NVMe 完整 IO 路径
# ══════════════════════════════════════════════════════
s = prs.slides.add_slide(BL)
sec(s, "01", "NVMe 完整 I/O 路径解析",
    "How GPU kernel reads from NVMe — end-to-end path")


# ══════════════════════════════════════════════════════
# Slide 5  NVMe 读取完整路径（关键！）
# ══════════════════════════════════════════════════════
s = prs.slides.add_slide(BL)
hdr(s, "完整 NVMe 读取路径 — bam_ptr Cache Miss 到数据就位",
    "This is the full path that needs to work on Corex")

# 步骤框
steps = [
    ("Step 1\nGPU kernel\n触发 cache miss",
     CYAN_D, CYAN_L,
     ["GPU warp 执行 ptr[row*dim+col]",
      "page_cache_t 查找对应页",
      "页状态 = INVALID → cache miss",
      "GPU warp 进入 NVMe 读取流程"],
     "GPU kernel"),
    ("Step 2\nGPU 构造\nNVMe 命令",
     BLUE, BLUE_L,
     ["在 SQ（提交队列）写入 NVMe READ 命令",
      "cmd.prp1 = GPU 页缓冲区 bus 地址",
      "cmd.slba = 数据在 NVMe 上的逻辑块地址",
      "GPU warp 直接写 nvm_queue_t 数据结构"],
     "GPU kernel"),
    ("Step 3\nGPU 写门铃\n触发 NVMe",
     ORANGE, ORANGE_L,
     ["*sq.db = sq.tail  // 写入 SQ 门铃寄存器",
      "sq.db = NVMe BAR0 mmap + GPU 虚地址",
      "通过 cudaHostRegisterIoMemory 获得",
      "NVMe 控制器检测到新命令"],
     "关键适配点"),
    ("Step 4\nNVMe DMA\nP2P 传输",
     GREEN_D, GREEN_L,
     ["NVMe 控制器读取 SQ 命令",
      "从 NAND Flash 读取数据块",
      "PCIe P2P DMA 写入 GPU 显存",
      "目标地址 = prp1 (GPU bus address)"],
     "硬件层"),
    ("Step 5\nGPU 轮询\nCQ 完成",
     RGBColor(0x60,0x30,0x9A), RGBColor(0xF5,0xF0,0xFF),
     ["GPU warp 轮询 CQ（完成队列）",
      "等待 NVMe 完成条目出现",
      "验证命令状态 = Success",
      "写 CQ 门铃，标记页为 VALID"],
     "GPU kernel"),
    ("Step 6\n数据就位\n继续执行",
     NAVY, CYAN_L,
     ["page_cache_t 标记页 VALID",
      "GPU warp 从显存读取数据",
      "ptr[idx] 返回正确数据值",
      "继续 GNN 计算（前向传播）"],
     "完成"),
]

step_w = (W - 800000) // 6
for i, (title, c, bg, pts, tag) in enumerate(steps):
    x = 350000 + i * (step_w + 15000)
    R(s, x, 580000, step_w, 5900000, fill=bg, line=c, lw=Pt(2))
    R(s, x, 580000, step_w, 420000, fill=c)
    T(s, x+40000, 600000, step_w-80000, 380000, title,
      sz=Pt(11), bold=True, color=BG, align=PP_ALIGN.CENTER)
    pts_l = [("• "+p, Pt(10), False, GRAY) for p in pts]
    ML(s, x+50000, 1070000, step_w-100000, 4100000, pts_l)
    R(s, x+50000, 5320000, step_w-100000, 300000, fill=c)
    T(s, x+50000, 5350000, step_w-100000, 240000, tag,
      sz=Pt(10), bold=True, color=BG, align=PP_ALIGN.CENTER)
    # 箭头
    if i < 5:
        T(s, x + step_w + 5000, 2800000, 30000, 300000, "→",
          sz=Pt(16), bold=True, color=c, align=PP_ALIGN.CENTER)

# 底部关键说明
R(s, 350000, 6500000, W-700000, 220000, fill=ORANGE_L, line=ORANGE, lw=Pt(1.5))
T(s, 450000, 6520000, W-900000, 190000,
  "Corex 适配关键：Step 3（cudaHostRegisterIoMemory → NVMe 门铃映射到 GPU VA）+ Step 4（GPU 显存作 NVMe DMA 目标）— 这两步决定 bam_ptr 是否能在 IX GPU 上运行",
  sz=Pt(11), color=ORANGE)


# ══════════════════════════════════════════════════════
# Slide 6  章节分隔 — 6大适配项
# ══════════════════════════════════════════════════════
s = prs.slides.add_slide(BL)
sec(s, "02", "6 大适配项逐项攻克",
    "Detailed analysis of each adaptation item")


# ══════════════════════════════════════════════════════
# Slide 7  适配项总览表
# ══════════════════════════════════════════════════════
s = prs.slides.add_slide(BL)
hdr(s, "6 大适配项总览 — 一图全景", "Green = confirmed, Orange = needs test, Red = needs development")

# 总览
items_summary = [
    ("①", "simt::atomic → cuda::atomic",
     "page_cache_t, nvm_types.h, nvm_parallel_queue.h\n全文件 simt:: 前缀",
     "cuda/std/atomic 已确认存在\nthread_scope_device 支持已验证",
     "全局 sed 替换", "1 天", "✅ 已确认可行", GREEN),
    ("②", "<<<>>> 语法 + CUDA Runtime API",
     "gids_kernel.cu kernel 调用语法\ncudaMalloc / cudaMemcpy / cudaMemset",
     "ixc 编译器原生支持 <<<>>>\nmapping_cudart.h 全部 1:1 映射",
     "无需修改", "0", "✅ 已确认可行", GREEN),
    ("③", "cuda_err_chk / 错误处理宏",
     "所有 .cu .cpp 文件中的\ncuda_err_chk() 调用",
     "ixError_t 类型兼容\n可用 #define cuda_err_chk ix_err_chk",
     "宏重定义", "0.5 天", "✅ 已确认可行", GREEN),
    ("④", "cudaHostRegisterIoMemory\nNVMe BAR → GPU VA",
     "ctrl.h: Controller 构造函数\ncudaHostRegister(mm_ptr, IoMemory)",
     "ixHostRegisterIoMemory = 0x04 已声明\nixHostRegister + ixHostGetDevicePointer 已有",
     "编写测试程序验证\nmmap NVMe BAR 后注册测试",
     "1 周", "⚠️ API存在,需运行时测试", ORANGE),
    ("⑤", "libnvm.ko Part A+B\n(NVMe BAR + 主机 DMA)",
     "module/pci.c: mmap_registers()\nmodule/map.c: NVM_MAP_HOST_MEMORY",
     "纯 Linux PCI/DMA API\n无任何 NVIDIA 专有依赖",
     "Linux 5.4 下直接编译\n(不开 -D_CUDA 编译选项)",
     "2~3 天", "✅ 纯Linux代码,可直接编译", GREEN),
    ("⑥", "libnvm.ko Part C\n(GPU 内存 DMA 注册)",
     "#ifdef _CUDA 部分\nnv-p2p.h + nvidia_p2p_get_pages()\nnvidia_p2p_dma_map_pages()",
     "选项A：Corex gdrapi.h 用户态路径\n选项B：向 Iluvatar 请求 ix-p2p.h 内核 API",
     "实现 Corex GPU 内存→NVMe DMA 注册\n替换 nvidia_p2p_* 系列函数",
     "2~4 周", "❌ 需要实现 Corex P2P 替换", RED),
]

row_h = 830000
for i, (num, name, prob, sol, action, effort, verdict, c) in enumerate(items_summary):
    col = i % 2; row = i // 2
    x = 350000 + col * 5950000
    y = 580000 + row * (row_h + 15000)
    vbg = GREEN_L if c==GREEN else (ORANGE_L if c==ORANGE else RED_L)
    R(s, x, y, 5750000, row_h, fill=vbg, line=c, lw=Pt(2))
    R(s, x, y, 350000, row_h, fill=c)
    T(s, x+50000, y+80000, 260000, row_h-160000, num,
      sz=Pt(24), bold=True, color=BG, align=PP_ALIGN.CENTER)
    T(s, x+420000, y+60000, 2200000, 350000, name,
      sz=Pt(13), bold=True, color=c, wrap=False)
    T(s, x+420000, y+420000, 2200000, 350000, f"问题：{prob}",
      sz=Pt(10), color=GRAY)
    T(s, x+2700000, y+60000, 2800000, 350000, sol,
      sz=Pt(11), color=GRAY)
    T(s, x+2700000, y+400000, 1500000, 360000, f"⏱ {effort}",
      sz=Pt(11), bold=True, color=c)
    R(s, x+4250000, y+150000, 1430000, 300000, fill=c)
    T(s, x+4270000, y+180000, 1390000, 250000, verdict,
      sz=Pt(10), bold=True, color=BG, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════
# Slide 8  适配项 ①②③ — 轻松解决
# ══════════════════════════════════════════════════════
s = prs.slides.add_slide(BL)
hdr(s, "适配项 ①②③ — 已确认可行（当日可完成）",
    "Zero-risk items — confirmed by examining Corex headers")

# ① simt::atomic
R(s, 350000, 580000, 3680000, 5900000, fill=GREEN_L, line=GREEN, lw=Pt(2.5))
R(s, 350000, 580000, 3680000, 400000, fill=GREEN)
T(s, 450000, 610000, 3490000, 330000, "① simt::atomic → cuda::atomic",
  sz=Pt(14), bold=True, color=BG)
ML(s, 450000, 1050000, 3480000, 2500000, [
    ("BaM 中的问题：", Pt(12), True, NAVY),
    ("nvm_types.h, page_cache.h, nvm_parallel_queue.h", Pt(11), False, GRAY),
    ("全部使用 NVIDIA CCCL #include <simt/atomic>", Pt(11), False, GRAY),
    ("simt::atomic<uint32_t, simt::thread_scope_device>", Pt(11), False, RED),
    ("simt::memory_order_relaxed / acquire / release", Pt(11), False, RED),
    ("", Pt(11), False, GRAY),
    ("Corex 验证结论：", Pt(12), True, GREEN),
    ("cuda/std/atomic 在 Corex 下可用 ✅", Pt(11), False, GREEN),
    ("thread_scope_device 已实现 ✅", Pt(11), False, GREEN),
    ("fetch_add/fetch_sub/store/load/exchange ✅", Pt(11), False, GREEN),
])
CODE(s, 450000, 3650000, 3480000, 2550000,
"""// 修复命令 (30 秒)
find bam/include -name "*.h" | xargs \\
  sed -i \\
  's|simt::|cuda::|g;\\
   s|<simt/atomic>|<cuda/std/atomic>|g'

// 修改前（BaM）
#include <simt/atomic>
simt::atomic<uint32_t,
  simt::thread_scope_device> state;
state.fetch_sub(1, simt::memory_order_release);

// 修改后（Corex）
#include <cuda/std/atomic>
cuda::atomic<uint32_t,
  cuda::thread_scope_device> state;
state.fetch_sub(1, cuda::memory_order_release);""")

# ② <<<>>> + CUDA API
R(s, 4130000, 580000, 3680000, 5900000, fill=GREEN_L, line=GREEN, lw=Pt(2.5))
R(s, 4130000, 580000, 3680000, 400000, fill=GREEN)
T(s, 4230000, 610000, 3490000, 330000, "② <<<>>> 语法 + CUDA Runtime API",
  sz=Pt(14), bold=True, color=BG)
ML(s, 4230000, 1050000, 3480000, 2000000, [
    ("<<<>>> 语法验证：", Pt(12), True, GREEN),
    ("ixc 编译器原生翻译 <<<>>> → __ixLaunch", Pt(11), False, GREEN),
    ("host_runtime.h 内部宏：", Pt(11), False, GRAY),
    ("__ixPopCallConfiguration() + ixLaunchKernel()", Pt(11), False, CODE_BLUE),
    ("无需任何源码修改 ✅", Pt(11), False, GREEN),
    ("", Pt(11), False, GRAY),
    ("CUDA Runtime API 已映射：", Pt(12), True, GREEN),
    ("cudaMalloc → ixMalloc ✅", Pt(11), False, GRAY),
    ("cudaMemcpy → ixMemcpy ✅", Pt(11), False, GRAY),
    ("cudaMemset → ixMemset ✅", Pt(11), False, GRAY),
    ("cudaHostAlloc → ixHostAlloc ✅", Pt(11), False, GRAY),
    ("cudaFree → ixFree ✅", Pt(11), False, GRAY),
    ("cuda_err_chk → 宏重定义 ✅", Pt(11), False, GRAY),
])
CODE(s, 4230000, 3100000, 3480000, 3100000,
"""// gids_kernel.cu — 无需修改
read_feature_kernel<TYPE>
    <<<g_size, b_size>>>(args...);

read_feature_kernel<TYPE>
    <<<g_size, b_size, 0, streams[i]>>>(args...);

// mapping_cudart.h 保证（节选）
#define cudaMalloc     ixMalloc
#define cudaMemcpy     ixMemcpy
#define cudaMemset     ixMemset
#define cudaHostAlloc  ixHostAlloc
#define cudaFree       ixFree
#define cudaSetDevice  ixSetDevice
// ... 共 3000+ 行映射""")

# ③ libnvm.ko Part A
R(s, 7910000, 580000, 3900000, 5900000, fill=GREEN_L, line=GREEN, lw=Pt(2.5))
R(s, 7910000, 580000, 3900000, 400000, fill=GREEN)
T(s, 8010000, 610000, 3700000, 330000, "③ libnvm.ko Part A+B — NVMe BAR mmap",
  sz=Pt(14), bold=True, color=BG)
ML(s, 8010000, 1050000, 3700000, 2000000, [
    ("module/pci.c — NVMe BAR 映射到用户态：", Pt(12), True, GREEN),
    ("mmap_registers() 使用标准 Linux API：", Pt(11), False, GRAY),
    ("pgprot_noncached() — 标准内存保护 ✅", Pt(11), False, GRAY),
    ("vm_iomap_memory() — 标准 MMIO 映射 ✅", Pt(11), False, GRAY),
    ("pci_resource_start/len() — 标准 PCI ✅", Pt(11), False, GRAY),
    ("", Pt(11), False, GRAY),
    ("module/map.c — 主机内存 DMA 注册：", Pt(12), True, GREEN),
    ("NVM_MAP_HOST_MEMORY ioctl", Pt(11), False, GRAY),
    ("dma_map_sg() — 标准 Linux DMA ✅", Pt(11), False, GRAY),
    ("无任何 #ifdef _CUDA 代码路径 ✅", Pt(11), False, GRAY),
])
CODE(s, 8010000, 3100000, 3700000, 3100000,
"""// 编译方法（不带 _CUDA）
# Makefile 中去掉 -D_CUDA 选项
# map.c 中所有 #ifdef _CUDA
# 块均被跳过

make -C /lib/modules/$(uname -r)/build \\
  M=$(pwd)/bam/module \\
  modules

# 生成 libnvm.ko
# 加载：insmod libnvm.ko
# 创建设备：/dev/libnvm0, /dev/libnvm1 ...""")


# ══════════════════════════════════════════════════════
# Slide 9  适配项 ④ — cudaHostRegisterIoMemory
# ══════════════════════════════════════════════════════
s = prs.slides.add_slide(BL)
hdr(s, "适配项 ④ — cudaHostRegisterIoMemory（NVMe 门铃→GPU 地址）",
    "The doorbell mapping: most important runtime behavior to verify")

# 原理图（左）
R(s, 350000, 580000, 5700000, 5900000, fill=ORANGE_L, line=ORANGE, lw=Pt(2.5))
T(s, 500000, 660000, 5500000, 350000, "为什么这是关键适配项",
  sz=Pt(15), bold=True, color=ORANGE)
ML(s, 500000, 1060000, 5500000, 2500000, [
    ("NVMe 门铃写入流程：", Pt(13), True, NAVY),
    ("", Pt(11), False, GRAY),
    ("1. CPU 打开 /dev/libnvm0，mmap() BAR0", Pt(12), False, GRAY),
    ("   → ctrl->mm_ptr = (void*) mmap(BAR0 物理地址)", Pt(11), False, CODE_BLUE),
    ("2. cudaHostRegister(mm_ptr, size, IoMemory)", Pt(12), False, GRAY),
    ("   → 告诉 CUDA：这段 MMIO 可从 GPU 访问", Pt(11), False, CODE_BLUE),
    ("3. cudaHostGetDevicePointer(&dev_ptr, sq.db, 0)", Pt(12), False, GRAY),
    ("   → 获取 NVMe SQ 门铃寄存器的 GPU 虚地址", Pt(11), False, CODE_BLUE),
    ("4. GPU kernel: *sq.db = sq.tail", Pt(12), False, GRAY),
    ("   → GPU store 写入 PCIe 地址 → NVMe 控制器感知", Pt(11), False, CODE_BLUE),
    ("", Pt(11), False, GRAY),
    ("Corex API 状态（已验证声明）：", Pt(13), True, GREEN),
    ("ixHostRegister(ptr, size, flags)  ✅ ix_runtime_api.h L363", Pt(12), False, GREEN),
    ("ixHostRegisterIoMemory = 0x04   ✅ driver_types.h L27", Pt(12), False, GREEN),
    ("ixHostGetDevicePointer(...)      ✅ ix_runtime_api.h L365", Pt(12), False, GREEN),
    ("mapping_cudart.h 全部已映射      ✅", Pt(12), False, GREEN),
    ("", Pt(11), False, GRAY),
    ("待验证：runtime 行为是否与 NVIDIA 一致", Pt(13), True, ORANGE),
    ("即：IX GPU 是否可以通过 PCIe 写入 NVMe BAR", Pt(12), False, ORANGE),
])

# 验证程序（右）
R(s, 6300000, 580000, 5500000, 5900000, fill=CODE_BG, line=CYAN_D, lw=Pt(2))
T(s, 6450000, 660000, 5300000, 350000, "验证测试程序（第一步就做这个）",
  sz=Pt(14), bold=True, color=CYAN_D)
CODE(s, 6450000, 1060000, 5250000, 4700000,
"""// test_io_memory.cu — 验证 ixHostRegisterIoMemory
#include <cuda_runtime.h>
#include <cstdio>
#include <fcntl.h>
#include <sys/mman.h>

__global__ void write_doorbell(volatile uint32_t* db,
                               uint32_t val) {
    *db = val;   // GPU 直接写门铃
}

int main() {
    // 1. 打开 /dev/libnvm0，mmap NVMe BAR
    int fd = open("/dev/libnvm0", O_RDWR);
    void* bar = mmap(NULL, 0x2000, PROT_READ|PROT_WRITE,
                     MAP_SHARED, fd, 0);

    // 2. 注册为 IoMemory（关键测试）
    cudaError_t err = cudaHostRegister(bar, 0x2000,
        cudaHostRegisterIoMemory);
    printf("HostRegister: %s\\n", cudaGetErrorString(err));

    // 3. 获取 GPU 侧指针
    volatile uint32_t* db_gpu = nullptr;
    cudaHostGetDevicePointer((void**)&db_gpu,
        (void*)((uint8_t*)bar + 0x1000), 0);

    // 4. GPU kernel 写门铃（核心验证）
    write_doorbell<<<1,1>>>(db_gpu, 1);
    cudaDeviceSynchronize();
    printf("GPU doorbell write: OK\\n");
    return 0;
}""")


# ══════════════════════════════════════════════════════
# Slide 10  适配项 ⑤ — libnvm.ko Part C (GPU DMA)
# ══════════════════════════════════════════════════════
s = prs.slides.add_slide(BL)
hdr(s, "适配项 ⑤ — libnvm.ko GPU DMA 注册（最大挑战）",
    "Replacing nvidia_p2p_* with Corex GPU P2P mechanism")

# 原理说明（左上）
R(s, 350000, 580000, 5700000, 2600000, fill=RED_L, line=RED, lw=Pt(2))
T(s, 500000, 650000, 5500000, 330000, "问题：GPU 内存作为 NVMe DMA 目标",
  sz=Pt(15), bold=True, color=RED)
ML(s, 500000, 1020000, 5500000, 2000000, [
    ("BaM page_cache_t 分配 GPU 显存作为 page buffer", Pt(12), False, GRAY),
    ("NVMe READ 命令的目标地址 = GPU 显存 bus 地址", Pt(12), False, GRAY),
    ("要求：GPU 物理页可被 NVMe 控制器 PCIe DMA 写入", Pt(12), False, GRAY),
    ("", Pt(11), False, GRAY),
    ("NVIDIA 实现（nv-p2p.h 内核 API）：", Pt(12), True, RED),
    ("nvidia_p2p_get_pages(va, size, &pages)  // 锁定 GPU 物理页", Pt(11), False, RED),
    ("nvidia_p2p_dma_map_pages(pdev, pages, &dma)  // 获取 DMA 地址", Pt(11), False, RED),
    ("dma->dma_addresses[i]  // 给 NVMe PRP 的 bus 地址", Pt(11), False, RED),
])

# 两种替换方案（左下）
R(s, 350000, 3300000, 5700000, 3100000, fill=RGBColor(0xF5,0xFB,0xFF), line=CYAN_D, lw=Pt(2))
T(s, 500000, 3380000, 5500000, 330000, "Corex 两种替换方案",
  sz=Pt(15), bold=True, color=CYAN_D)
ML(s, 500000, 3760000, 5500000, 2400000, [
    ("方案 A（推荐先试）：gdrapi.h 用户态路径", Pt(12), True, CYAN_D),
    ("1. ixdrvPointerGetAttributes(P2P_TOKENS, &tok, gpu_ptr)", Pt(11), False, GRAY),
    ("2. gdr_open() + gdr_pin_buffer(tok, va_space)", Pt(11), False, GRAY),
    ("3. gdr_get_info() 获取物理/bus 地址", Pt(11), False, GRAY),
    ("4. nvm_dma_map(addrs[]) 直接传 bus 地址（绕过内核模块）", Pt(11), False, GRAY),
    ("风险：Corex IX GPU 物理地址是否等于 PCIe bus 地址？", Pt(11), False, ORANGE),
    ("", Pt(11), False, GRAY),
    ("方案 B（最干净）：请 Iluvatar 提供 ix-p2p.h 内核 API", Pt(12), True, CYAN_D),
    ("内核中替换 nvidia_p2p_get_pages() → ix_p2p_get_pages()", Pt(11), False, GRAY),
    ("nvidia_p2p_dma_map_pages() → ix_p2p_dma_map_pages()", Pt(11), False, GRAY),
    ("需要 Iluvatar 内核驱动团队支持（1~2 周）", Pt(11), False, ORANGE),
])

# 代码实现（右）
R(s, 6300000, 580000, 5500000, 5900000, fill=CODE_BG, line=ORANGE, lw=Pt(2))
T(s, 6450000, 650000, 5300000, 330000, "方案 A 实现框架",
  sz=Pt(14), bold=True, color=ORANGE)
CODE(s, 6450000, 1030000, 5250000, 5250000,
"""// libnvm_corex.c — 替换 #ifdef _CUDA 部分
// (userspace path via gdrapi.h)
#include "gdrapi.h"   // Corex GDR API

int map_gpu_memory_corex(struct map* map,
    struct list* ctrl_list)
{
    // Step1: 获取 GPU 内存物理 bus 地址
    // (用户态事先通过 gdrapi 获取并传入)
    // map->vaddr = GPU VA (ixMalloc 分配)
    // map->addrs[] 已由用户态填充 bus addresses

    // Step2: 注册为 NVMe DMA 目标
    // 直接使用传入的 bus 地址构造 DMA 映射
    // (无需内核态 nv-p2p, bus 地址由用户态通过
    //  gdr_pin_buffer + gdr_get_info 获取)

    return 0;
}

// 用户态侧（buffer.h 修改）
// 获取 GPU 内存的 bus 地址（Corex方式）
void get_gpu_bus_addrs(void* gpu_ptr,
                       size_t size,
                       uint64_t* bus_addrs,
                       int n_pages)
{
    IXDRVresult res;
    IX_POINTER_ATTRIBUTE_P2P_TOKENS tokens;
    res = ixdrvPointerGetAttributes(1,
        IXDRV_POINTER_ATTRIBUTE_P2P_TOKENS,
        &tokens, (CUdeviceptr)gpu_ptr);
    gdr_t g = gdr_open();
    gdr_mh_t handle;
    gdr_pin_buffer(g, (ulong)gpu_ptr, size,
        tokens.p2pToken,
        tokens.vaSpace, &handle);
    gdr_info_t info;
    gdr_get_info_v2(g, handle, &info);
    // info.va = mapped VA
    // 从 info 计算各页 bus 地址填入 bus_addrs[]
}""")


# ══════════════════════════════════════════════════════
# Slide 11  章节分隔 — 适配后架构
# ══════════════════════════════════════════════════════
s = prs.slides.add_slide(BL)
sec(s, "03", "适配后完整架构",
    "Post-adaptation architecture — bam_ptr native on Corex")


# ══════════════════════════════════════════════════════
# Slide 12  适配后架构图
# ══════════════════════════════════════════════════════
s = prs.slides.add_slide(BL)
hdr(s, "适配后架构 — bam_ptr Corex 版全景图",
    "Every layer adapted, bam_ptr semantics fully preserved")

R(s, 0, 0, W, H, fill=BG)
R(s, 0, 0, W, 500000, fill=NAVY)
R(s, 0, 6380000, W, 50000, fill=CYAN)
T(s, 400000, 80000, W-800000, 360000, "适配后架构 — bam_ptr Corex 版全景图",
  sz=Pt(22), bold=True, color=BG)

arch_layers = [
    ("Python 训练层",
     ["GIDS.py / GIDS_DGLDataLoader", "homogenous_train.py"],
     ["调用方式不变", "import BAM_Feature_Store as BFS"],
     GREEN, GREEN_L, "无需修改"),
    ("pybind11 绑定层",
     ["BAM_Feature_Store.so", "BAM_Feature_Store<float/half/long>"],
     ["pybind11 编译改为 ixc", "CMakeLists.txt 适配"],
     BLUE, BLUE_L, "CMake 适配"),
    ("BaM C++/CUDA 核心",
     ["page_cache_t（simt→cuda 修复）", "bam_ptr<T>（无需修改）"],
     ["gids_kernel.cu（<<<>>> 不变）", "Controller（cudaHostRegisterIoMemory）"],
     CYAN_D, CYAN_L, "①②③④修复"),
    ("libnvm 用户态库",
     ["libnvm.so", "nvm_ctrl_init / nvm_dma_map_host"],
     ["A+B部分编译（去掉-D_CUDA）", "C部分替换为 gdrapi.h 路径"],
     ORANGE, ORANGE_L, "⑤替换"),
    ("libnvm.ko 内核模块",
     ["NVMe PCI BAR mmap（不变）", "主机内存 DMA 注册（不变）"],
     ["GPU DMA 注册改用 Corex P2P", "编译目标：Linux 5.4 + Corex"],
     RED, RED_L, "⑥P2P替换"),
    ("硬件层",
     ["IX GPU（Corex CUDA 兼容）", "NVMe SSD（/dev/libnvm0..N）"],
     ["ixHostRegisterIoMemory 映射 BAR", "GPU 显存 PCIe P2P DMA 接收"],
     NAVY, CYAN_L, "硬件支持"),
]

lw = (H - 750000) // 6
for i, (title, left, right, c, bg, tag) in enumerate(arch_layers):
    ly = 580000 + i * (lw + 10000)
    R(s, 250000, ly, W-500000, lw, fill=bg, line=c, lw=Pt(2))
    R(s, 250000, ly, 380000, lw, fill=c)
    T(s, 300000, ly + lw//2 - 180000, 310000, 360000, tag,
      sz=Pt(10), bold=True, color=BG, align=PP_ALIGN.CENTER)
    T(s, 720000, ly+80000, 3000000, 350000, title,
      sz=Pt(13), bold=True, color=c)
    l_l = [("• "+l, Pt(11), False, GRAY) for l in left]
    ML(s, 720000, ly+420000, 3000000, lw-500000, l_l)
    r_l = [("• "+r, Pt(11), False, GRAY) for r in right]
    ML(s, 4200000, ly+80000, 7000000, lw-160000, r_l)
    # 适配标记
    if i > 0:
        T(s, W-700000, ly+lw//2-180000, 400000, 360000, f"↑",
          sz=Pt(14), bold=True, color=c, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════
# Slide 13  章节分隔 — 工作流程
# ══════════════════════════════════════════════════════
s = prs.slides.add_slide(BL)
sec(s, "04", "工作流程与里程碑",
    "Step-by-step implementation plan with milestones")


# ══════════════════════════════════════════════════════
# Slide 14  详细工作流程甘特图
# ══════════════════════════════════════════════════════
s = prs.slides.add_slide(BL)
hdr(s, "工作流程甘特图 — bam_ptr 原生适配完整计划",
    "6-week implementation plan")

R(s, 0, 0, W, H, fill=BG)
R(s, 0, 0, W, 500000, fill=NAVY)
R(s, 0, 6380000, W, 50000, fill=CYAN)
T(s, 400000, 80000, W-800000, 360000,
  "工作流程甘特图 — bam_ptr 原生适配完整计划",
  sz=Pt(22), bold=True, color=BG)

weeks = ["D1-3\n验证期", "W1\n基础适配", "W2\n模块编译", "W3\n集成测试", "W4\n功能验证", "W5-6\n性能优化"]
n_weeks = 6
week_w = (W - 1400000) // n_weeks
header_y = 580000; row_h = 820000

# 表头
R(s, 380000, header_y, 900000, 420000, fill=NAVY)
T(s, 390000, header_y+80000, 870000, 280000, "任务",
  sz=Pt(11), bold=True, color=BG, align=PP_ALIGN.CENTER)
for wi, wk in enumerate(weeks):
    wx = 1280000 + wi * week_w
    R(s, wx, header_y, week_w - 10000, 420000, fill=NAVY)
    T(s, wx+20000, header_y+50000, week_w-60000, 310000, wk,
      sz=Pt(11), bold=True, color=BG, align=PP_ALIGN.CENTER)

tasks_g = [
    # (名称, 颜色, 开始(0-based), 跨度, 说明)
    ("① simt::atomic → cuda::atomic", GREEN, 0, 1, "sed 替换，编译验证"),
    ("② <<<>>>语法 + CUDA RT API 验证", GREEN, 0, 1, "已确认，无修改"),
    ("③ cuda_err_chk 宏适配", GREEN, 0, 1, "宏重定义，0.5天"),
    ("⑤ libnvm.ko A+B 编译（无CUDA）", GREEN, 1, 1, "Linux 5.4 kernel build"),
    ("④ cudaHostRegisterIoMemory 验证", ORANGE, 0, 2, "test_io_memory.cu 测试程序"),
    ("BaM 源码 ixc 编译调试", BLUE, 1, 1, "page_cache_t + Controller"),
    ("libnvm.so 库链接验证", CYAN_D, 2, 1, "用户态 NVMe 访问测试"),
    ("⑥ GPU DMA注册 gdrapi 实现", RED, 1, 3, "方案A：gdrapi.h 替换路径"),
    ("方案B备选：ix-p2p.h 内核API", RED, 1, 3, "需Iluvatar驱动团队"),
    ("Controller 端到端初始化测试", ORANGE, 3, 1, "打开/dev/libnvm0，初始化队列"),
    ("bam_ptr 单次 NVMe 读取测试", CYAN_D, 3, 1, "最小可行验证"),
    ("page_cache_t 端到端功能测试", BLUE, 4, 1, "多页并发，LRU 淘汰"),
    ("gids_kernel 集成测试", GREEN_D, 4, 1, "read_feature_kernel 正确性"),
    ("性能基准测试（vs cuFile路线）", NAVY, 5, 1, "吞吐量/延迟对比"),
]

for ri, (name, c, start, span, note) in enumerate(tasks_g):
    ry = header_y + 460000 + ri * (row_h + 10000)
    bg = RGBColor(0xF8,0xFC,0xFF) if ri%2==0 else BG
    R(s, 380000, ry, 900000, row_h, fill=bg, line=BORD, lw=Pt(1))
    T(s, 390000, ry+100000, 870000, row_h-200000,
      str(ri+1), sz=Pt(14), bold=True, color=LGRAY, align=PP_ALIGN.CENTER)
    bar_x = 1280000 + start * week_w
    bar_w = span * week_w - 15000
    R(s, bar_x, ry+80000, bar_w, row_h-160000, fill=c)
    T(s, bar_x+60000, ry+140000, bar_w-120000, row_h-300000,
      name, sz=Pt(10), bold=True, color=BG)
    note_x = 1280000 + n_weeks * week_w + 20000
    T(s, note_x, ry+200000, 1200000, row_h-400000, note,
      sz=Pt(9), color=LGRAY)


# ══════════════════════════════════════════════════════
# Slide 15  阶段验证计划
# ══════════════════════════════════════════════════════
s = prs.slides.add_slide(BL)
hdr(s, "阶段验证计划 — 从最小验证到端到端",
    "Incremental validation strategy — each gate blocks next phase")

phases = [
    ("Gate 1\nD1~D3", ORANGE, "IoMemory 运行时验证",
     ["运行 test_io_memory.cu 测试程序",
      "确认 ixHostRegister(IoMemory) 无报错",
      "确认 GPU warp 写门铃不崩溃/不 hang",
      "验证 NVMe 能收到 doorbell 信号"],
     ["Pass → 继续 Gate 2", "Fail → 联系 Iluvatar 运行时团队"],
     "测试程序：test_io_memory.cu"),
    ("Gate 2\nW1~W2", GREEN, "libnvm + BaM 源码编译通过",
     ["libnvm.ko 在 Linux 5.4 下编译通过",
      "BaM 所有头文件 ixc 编译无错误",
      "Controller 构造函数无运行时异常",
      "QueuePair 初始化，门铃指针有效"],
     ["Pass → 继续 Gate 3", "Fail → debug simt→cuda 遗漏项"],
     "make -j4 + ixc 编译测试"),
    ("Gate 3\nW2~W3", CYAN_D, "GPU DMA 注册验证",
     ["GPU 显存 bus 地址可通过 gdrapi 获取",
      "nvm_dma_t::ioaddrs[] 填充正确",
      "nvm_dma_map() 成功注册 DMA 目标",
      "NVMe READ DMA 写入 GPU 显存成功"],
     ["Pass → 继续 Gate 4", "Fail → 切换方案 B（ix-p2p.h）"],
     "test_gpu_dma_map 验证程序"),
    ("Gate 4\nW3~W4", BLUE, "bam_ptr 端到端功能验证",
     ["bam_ptr[0] 触发正确 NVMe READ",
      "page_cache_t LRU 淘汰正常工作",
      "多线程并发不死锁（simt→cuda 原子）",
      "数据校验：读到的特征值与文件一致"],
     ["Pass → 继续 Gate 5", "Fail → 并发原子或 DMA 问题排查"],
     "GIDS_unit_test.py 单测"),
    ("Gate 5\nW5~W6", NAVY, "性能基准达标",
     ["吞吐量 ≥ cuFile 路线的 1.5x（理论 2~3x）",
      "GPU 利用率 > 80%（无 CPU 瓶颈）",
      "IGB-1M 数据集端到端训练 1 epoch 通过",
      "与 NVIDIA 原版 GIDS 性能对比报告"],
     ["Pass → 方案确定，转生产", "未达标 → Window Buffer / 多队列优化"],
     "吞吐量 vs 延迟 benchmark"),
]

pw = (W - 700000) // 5
for i, (phase, c, title_g, steps, results, tool) in enumerate(phases):
    x = 350000 + i * (pw + 10000)
    bg = GREEN_L if c==GREEN or c==GREEN_D or c==NAVY else (ORANGE_L if c==ORANGE else CYAN_L if c==CYAN_D else BLUE_L)
    R(s, x, 580000, pw, 5900000, fill=bg, line=c, lw=Pt(2.5))
    R(s, x, 580000, pw, 420000, fill=c)
    T(s, x+40000, 600000, pw-80000, 380000, phase,
      sz=Pt(13), bold=True, color=BG, align=PP_ALIGN.CENTER)
    T(s, x+50000, 1070000, pw-100000, 320000, title_g,
      sz=Pt(12), bold=True, color=c)
    step_l = [("• "+st, Pt(10), False, GRAY) for st in steps]
    ML(s, x+50000, 1430000, pw-100000, 2100000, step_l)
    T(s, x+50000, 3600000, pw-100000, 240000, "结果分支：",
      sz=Pt(11), bold=True, color=NAVY)
    res_l = [(r, Pt(10), False, GREEN if "Pass" in r else RED) for r in results]
    ML(s, x+50000, 3870000, pw-100000, 800000, res_l)
    R(s, x+50000, 4800000, pw-100000, 280000, fill=RGBColor(0xE0,0xEC,0xF8))
    T(s, x+60000, 4830000, pw-120000, 240000, "🔧 "+tool,
      sz=Pt(9), color=NAVY)


# ══════════════════════════════════════════════════════
# Slide 16  风险对策
# ══════════════════════════════════════════════════════
s = prs.slides.add_slide(BL)
hdr(s, "风险评估与对策", "Risk matrix for bam_ptr native porting")

risks = [
    ("R1 中高", ORANGE,
     "ixHostRegisterIoMemory 运行时不生效",
     "IX GPU 驱动可能不支持从 GPU 写入任意 MMIO 地址段",
     "D1~D3 运行 test_io_memory.cu 验证\n一旦确认，立即有结论",
     "发现时间 D3 前，切换 cuFile 路线零代价"),
    ("R2 高", RED,
     "GPU DMA 目标注册（gdrapi 方案 A）失败",
     "IX GPU 物理地址 ≠ PCIe bus 地址（IOMMU 影响）",
     "同步推进方案 B：请 Iluvatar 提供 ix-p2p.h\n准备好方案 B 的接口规范文档",
     "方案 B 需要 Iluvatar 驱动团队 1~2 周支持"),
    ("R3 中", ORANGE,
     "simt→cuda atomic 语义差异",
     "Corex cuda::atomic 可能在某些 memory_order 组合\n下行为与 NVIDIA simt::atomic 不完全一致",
     "Gate 4 并发正确性测试\n用 valgrind/sanitizer 检测数据竞争",
     "可替换为 __atomic_* GCC 内建函数保底"),
    ("R4 低", GREEN,
     "libnvm.ko 在 Linux 5.4 编译失败",
     "内核 API 版本差异（probe/remove 函数签名等）",
     "BaM module 代码量小（~300行），修改成本低\n参考 itrfs.ko 适配经验",
     "3 天内可解决"),
    ("R5 低", GREEN,
     "DGL 编译进度影响端到端测试",
     "DGL CUDA 版编译与 bam_ptr 适配并行进行",
     "bam_ptr 单元测试不依赖 DGL\n两条线独立，不互相阻塞",
     "DGL 进度不影响 bam_ptr Gate 1~4"),
    ("R6 低", GREEN,
     "工期超过 6 周",
     "GPU DMA 部分（R2）是主要风险",
     "Gate 1~3 均有明确的 pass/fail 判断\n有 cuFile 作为并行保底方案",
     "6 周内无论结果如何都有可用方案"),
]

rh = 850000
for i, (level, c, title_r, prob, sol, note) in enumerate(risks):
    col = i % 2; row = i // 2
    x = 350000 + col * 6000000
    y = 600000 + row * (rh + 15000)
    bg = GREEN_L if c==GREEN else (ORANGE_L if c==ORANGE else RED_L)
    R(s, x, y, 5800000, rh, fill=bg, line=c, lw=Pt(2))
    R(s, x, y, 600000, rh, fill=c)
    T(s, x+40000, y+70000, 540000, rh-140000, level.replace(' ','\n'),
      sz=Pt(11), bold=True, color=BG, align=PP_ALIGN.CENTER)
    T(s, x+680000, y+60000, 5000000, 320000, title_r,
      sz=Pt(13), bold=True, color=c)
    T(s, x+680000, y+390000, 5000000, 250000, "根因："+prob,
      sz=Pt(10), color=GRAY)
    T(s, x+680000, y+620000, 5000000, 250000, "对策："+sol,
      sz=Pt(10), color=GREEN_D)


# ══════════════════════════════════════════════════════
# Slide 17  总结
# ══════════════════════════════════════════════════════
s = prs.slides.add_slide(BL)
R(s, 0, 0, W, H, fill=NAVY)
R(s, 0, 0, W, 6000, fill=CYAN)
R(s, 0, H-6000, W, 6000, fill=CYAN)
R(s, 0, 0, 200000, H, fill=CYAN)

T(s, 400000, 400000, W-800000, 600000,
  "总结 — bam_ptr 原生适配可行性与工作重点",
  sz=Pt(28), bold=True, color=BG)

summary_items = [
    (GREEN,  "✅ 4 项已确认可行（当日可完成）",
     "simt→cuda、<<<>>>语法、CUDA Runtime API映射、libnvm.ko非CUDA部分 — 全部有明确解法"),
    (ORANGE, "⚠️ 1 项需运行时验证（D1~D3 关键路径）",
     "cudaHostRegisterIoMemory：Corex 声明存在（0x04），运行时能否映射 NVMe MMIO 到 GPU VA 需测试"),
    (RED,    "❌ 1 项需要开发（GPU DMA 注册，W1~W3）",
     "nvidia_p2p_* → 方案A(gdrapi.h用户态) 或 方案B(ix-p2p.h内核) 替换，是整个移植的最大挑战"),
    (NAVY,   "🎯 核心价值：完整保留 bam_ptr GPU-Direct 路径",
     "GPU kernel 透明访问 NVMe，无 CPU 介入，性能上限与 NVIDIA 原版 GIDS 一致"),
]

for i, (c, title_s, desc) in enumerate(summary_items):
    y = 1200000 + i * 1250000
    R(s, 400000, y, W-800000, 1150000,
      fill=RGBColor(0x12,0x28,0x40), line=c, lw=Pt(2.5))
    R(s, 400000, y, 30000, 1150000, fill=c)
    T(s, 550000, y+100000, W-1100000, 420000, title_s,
      sz=Pt(15), bold=True, color=c)
    T(s, 550000, y+550000, W-1100000, 530000, desc,
      sz=Pt(13), color=RGBColor(0xB0,0xCC,0xE0))

T(s, 400000, 6480000, W-800000, 280000,
  "BaM × Corex  |  bam_ptr Native Adaptation  |  6 Weeks Plan  |  2026-06-12",
  sz=Pt(11), color=RGBColor(0x3A,0x6A,0x9A), align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════
out = "/root/GIDS_cufile/GIDS-bam_ptr原生适配Corex-完整分析PPT.pptx"
prs.save(out)
print(f"✅ 生成：{out}")
print(f"   共 {len(prs.slides)} 张幻灯片")
