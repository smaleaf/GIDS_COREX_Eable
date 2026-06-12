#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
GIDS 架构分析 PPT 生成脚本
风格参考：CUDA-GPU-Programming-TMA-PPT.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ─── 品牌配色（对齐参考 PPT）─────────────────────────────
BG_WHITE       = RGBColor(0xFF, 0xFF, 0xFF)   # 背景白
ACCENT_CYAN    = RGBColor(0x51, 0xB9, 0xD6)   # 主强调色 青蓝
TITLE_NAVY     = RGBColor(0x44, 0x54, 0x6A)   # 标题深灰蓝
SUBTITLE_BLUE  = RGBColor(0x2E, 0x86, 0xC1)   # 副标题蓝
BODY_GRAY      = RGBColor(0x60, 0x6A, 0x78)   # 正文灰
SECTION_BG     = RGBColor(0x1A, 0x3A, 0x5C)   # 章节页深蓝背景
HIGHLIGHT_YEL  = RGBColor(0xFF, 0xA5, 0x00)   # 警告橙
SUCCESS_GREEN  = RGBColor(0x27, 0xAE, 0x60)   # 成功绿
ERROR_RED      = RGBColor(0xE7, 0x4C, 0x3C)   # 错误红
LIGHT_BG       = RGBColor(0xF0, 0xF7, 0xFF)   # 浅蓝背景
BORDER_COLOR   = RGBColor(0xD5, 0xE8, 0xF5)   # 边框浅蓝
CODE_BG        = RGBColor(0xF5, 0xF5, 0xF5)   # 代码块背景

SLIDE_W = 12191365
SLIDE_H = 6858000


def rgb(r, g, b):
    return RGBColor(r, g, b)


def add_rect(slide, x, y, w, h, fill_color=None, line_color=None, line_width=Pt(1)):
    from pptx.util import Emu
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Emu(x), Emu(y), Emu(w), Emu(h)
    )
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()

    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, x, y, w, h, text, font_size=Pt(14), bold=False,
                color=BODY_GRAY, align=PP_ALIGN.LEFT, font_name="微软雅黑",
                wrap=True):
    from pptx.util import Emu
    txBox = slide.shapes.add_textbox(Emu(x), Emu(y), Emu(w), Emu(h))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    run.font.name = font_name
    return txBox


def add_multiline_textbox(slide, x, y, w, h, lines, font_size=Pt(13),
                          color=BODY_GRAY, font_name="微软雅黑",
                          line_spacing=None, align=PP_ALIGN.LEFT):
    from pptx.util import Emu
    txBox = slide.shapes.add_textbox(Emu(x), Emu(y), Emu(w), Emu(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        if isinstance(line, tuple):
            text, size, bold, col = line
        else:
            text, size, bold, col = line, font_size, False, color
        run = p.add_run()
        run.text = text
        run.font.size = size
        run.font.bold = bold
        run.font.color.rgb = col
        run.font.name = font_name
    return txBox


def add_title_bar(slide, title, subtitle=None):
    """顶部标题栏（带青蓝色条）"""
    # 白色背景
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill_color=BG_WHITE)
    # 顶部青蓝横条
    add_rect(slide, 0, 0, SLIDE_W, 457200, fill_color=ACCENT_CYAN)
    # 底部细线
    add_rect(slide, 0, 6400800, SLIDE_W, 54864, fill_color=ACCENT_CYAN)
    # 标题文字
    add_textbox(slide, 457200, 80000, SLIDE_W - 914400, 320000,
                title, font_size=Pt(24), bold=True,
                color=BG_WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_textbox(slide, 457200, 6430000, SLIDE_W - 914400, 350000,
                    subtitle, font_size=Pt(11), bold=False,
                    color=BODY_GRAY, align=PP_ALIGN.RIGHT)


def add_section_divider(slide, number, title, subtitle=""):
    """章节分隔页"""
    # 深蓝背景
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill_color=SECTION_BG)
    # 青蓝左竖条
    add_rect(slide, 0, 0, 220000, SLIDE_H, fill_color=ACCENT_CYAN)
    # 章节编号
    add_textbox(slide, 500000, 2000000, 3000000, 1200000,
                number, font_size=Pt(80), bold=True,
                color=ACCENT_CYAN, align=PP_ALIGN.LEFT)
    # 标题
    add_textbox(slide, 500000, 3100000, SLIDE_W - 1000000, 800000,
                title, font_size=Pt(36), bold=True,
                color=BG_WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_textbox(slide, 500000, 3900000, SLIDE_W - 1000000, 500000,
                    subtitle, font_size=Pt(18), bold=False,
                    color=ACCENT_CYAN, align=PP_ALIGN.LEFT)


def add_bullet_box(slide, x, y, w, h, title, bullets, bg=LIGHT_BG,
                   border=BORDER_COLOR, title_color=SUBTITLE_BLUE,
                   bullet_color=BODY_GRAY):
    """带标题的项目符号框"""
    add_rect(slide, x, y, w, h, fill_color=bg, line_color=border, line_width=Pt(1.5))
    # 标题
    add_textbox(slide, x + 150000, y + 80000, w - 300000, 280000,
                title, font_size=Pt(14), bold=True, color=title_color)
    # 项目符号
    lines = [("• " + b, Pt(12), False, bullet_color) for b in bullets]
    add_multiline_textbox(slide, x + 150000, y + 350000, w - 300000,
                          h - 450000, lines)


def add_status_badge(slide, x, y, text, status="ok"):
    """状态徽章"""
    if status == "ok":
        bg, fg = rgb(0x27, 0xAE, 0x60), BG_WHITE
        prefix = "✅ "
    elif status == "warn":
        bg, fg = rgb(0xFF, 0xA5, 0x00), BG_WHITE
        prefix = "⚠️ "
    elif status == "err":
        bg, fg = rgb(0xE7, 0x4C, 0x3C), BG_WHITE
        prefix = "❌ "
    else:
        bg, fg = ACCENT_CYAN, BG_WHITE
        prefix = "◈ "
    w, h = 2200000, 280000
    add_rect(slide, x, y, w, h, fill_color=bg)
    add_textbox(slide, x + 80000, y + 30000, w - 160000, h - 60000,
                prefix + text, font_size=Pt(11), bold=True, color=fg,
                align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════
# 创建 PPT
# ════════════════════════════════════════════════════════════════
prs = Presentation()
prs.slide_width = Emu(SLIDE_W)
prs.slide_height = Emu(SLIDE_H)

blank_layout = prs.slide_layouts[6]  # Blank


# ── Slide 1: 封面 ──────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)

# 白色全背景
add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill_color=BG_WHITE)

# 顶部装饰条
add_rect(slide, 0, 0, SLIDE_W, 100000, fill_color=ACCENT_CYAN)
# 中央青蓝分隔条
add_rect(slide, 0, 2600000, SLIDE_W, 60000, fill_color=ACCENT_CYAN)
# 底部装饰条
add_rect(slide, 0, SLIDE_H - 100000, SLIDE_W, 100000, fill_color=SECTION_BG)
# 左侧深色竖条
add_rect(slide, 0, 0, 180000, SLIDE_H, fill_color=SECTION_BG)

# 主标题
add_textbox(slide, 500000, 800000, SLIDE_W - 1000000, 900000,
            "GIDS GPU Direct Storage", font_size=Pt(40), bold=True,
            color=TITLE_NAVY, align=PP_ALIGN.CENTER)
add_textbox(slide, 500000, 1700000, SLIDE_W - 1000000, 700000,
            "架构分析 · Corex 适配方案 · 后期工作规划", font_size=Pt(22),
            bold=True, color=SUBTITLE_BLUE, align=PP_ALIGN.CENTER)

# 副信息
add_textbox(slide, 500000, 2850000, SLIDE_W - 1000000, 400000,
            "GPU-Initiated Direct Storage Accesses for GNN Training",
            font_size=Pt(16), bold=False, color=BODY_GRAY, align=PP_ALIGN.CENTER)

# 信息列
info_x = [800000, 3800000, 6800000, 9800000]
info_labels = ["项目来源", "适配目标", "当前状态", "日期"]
info_values = ["ZaidQureshi/bam", "Iluvatar Corex 4.5.0", "核心代码移植完成", "2026-06-12"]
for ix, (label, val) in enumerate(zip(info_labels, info_values)):
    x = info_x[ix]
    add_rect(slide, x, 3600000, 2300000, 400000, fill_color=LIGHT_BG, line_color=BORDER_COLOR)
    add_textbox(slide, x + 80000, 3630000, 2150000, 160000, label,
                font_size=Pt(10), bold=False, color=BODY_GRAY)
    add_textbox(slide, x + 80000, 3790000, 2150000, 180000, val,
                font_size=Pt(12), bold=True, color=SUBTITLE_BLUE)

# 底部
add_textbox(slide, 500000, 4400000, SLIDE_W - 1000000, 300000,
            "Iluvatar CoreX Platform  |  GNN Training Acceleration  |  GPU Direct Storage",
            font_size=Pt(12), bold=False, color=BODY_GRAY, align=PP_ALIGN.CENTER)


# ── Slide 2: 目录 ──────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
add_title_bar(slide, "目录", "Table of Contents")

toc_items = [
    ("01", "GIDS 项目概述", "核心思想、关键特性与系统定位"),
    ("02", "系统架构总览", "四层架构分解：Python → pybind11 → C++/CUDA → BaM/Hardware"),
    ("03", "核心组件深度解析", "BaM Framework、GPU Kernels、BAM_Feature_Store"),
    ("04", "三大优化策略", "Window Buffering / Storage Accumulator / CPU Buffer"),
    ("05", "Corex 适配挑战分析", "硬件差异、软件依赖、移植风险评估"),
    ("06", "适配方案与当前进展", "cuFile 替代 BaM、IXFeatureStore 设计"),
    ("07", "关键适配工作清单", "DGL 编译、cooperative_groups 补丁、cuFile 集成"),
    ("08", "后期工作流程与里程碑", "端到端验证路径、阶段目标"),
]

for i, (num, title, sub) in enumerate(toc_items):
    col = i % 2
    row = i // 2
    x = 457200 + col * 5700000
    y = 600000 + row * 1350000
    w = 5400000
    h = 1200000
    add_rect(slide, x, y, w, h, fill_color=LIGHT_BG, line_color=BORDER_COLOR, line_width=Pt(1.5))
    add_textbox(slide, x + 80000, y + 80000, 800000, 500000,
                num, font_size=Pt(28), bold=True, color=ACCENT_CYAN, align=PP_ALIGN.LEFT)
    add_textbox(slide, x + 850000, y + 80000, w - 950000, 380000,
                title, font_size=Pt(15), bold=True, color=TITLE_NAVY)
    add_textbox(slide, x + 850000, y + 460000, w - 950000, 600000,
                sub, font_size=Pt(11), bold=False, color=BODY_GRAY)


# ── Slide 3: 章节分隔 - 项目概述 ───────────────────────────────
slide = prs.slides.add_slide(blank_layout)
add_section_divider(slide, "01", "GIDS 项目概述",
                    "GPU-Initiated Direct Storage Accesses")


# ── Slide 4: 项目概述 ──────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
add_title_bar(slide, "GIDS 项目概述", "GPU Direct Storage for GNN Training")

# 左侧：核心思想
add_rect(slide, 457200, 600000, 5000000, 5600000, fill_color=LIGHT_BG, line_color=BORDER_COLOR)
add_textbox(slide, 600000, 680000, 4700000, 330000,
            "核心思想", font_size=Pt(16), bold=True, color=SUBTITLE_BLUE)
add_rect(slide, 600000, 1000000, 4700000, 4, fill_color=ACCENT_CYAN)

core_idea_lines = [
    ("GPU Direct Storage（GDS）", Pt(13), True, TITLE_NAVY),
    ("让 GPU kernel 直接从 NVMe SSD 读取图节点特征数据", Pt(12), False, BODY_GRAY),
    ("", Pt(12), False, BODY_GRAY),
    ("传统痛点：", Pt(13), True, ERROR_RED),
    ("CPU 瓶颈 → 数据流水线成为 GNN 训练的主要瓶颈", Pt(12), False, BODY_GRAY),
    ("PCIe 两次传输：SSD→RAM→GPU（带宽浪费）", Pt(12), False, BODY_GRAY),
    ("CPU 内存容量限制（TB 级图特征无法全量缓存）", Pt(12), False, BODY_GRAY),
    ("", Pt(12), False, BODY_GRAY),
    ("GIDS 解决方案：", Pt(13), True, SUCCESS_GREEN),
    ("绕过 CPU，GPU 直接访问 NVMe（P2P DMA）", Pt(12), False, BODY_GRAY),
    ("三级缓存体系：GPU显存 + CPU pinned + NVMe", Pt(12), False, BODY_GRAY),
    ("异步预取 + 批量合并，消除 SSD 访问延迟", Pt(12), False, BODY_GRAY),
]
add_multiline_textbox(slide, 600000, 1050000, 4700000, 4900000, core_idea_lines)

# 右侧：关键特性
add_rect(slide, 5800000, 600000, 5900000, 2650000, fill_color=LIGHT_BG, line_color=BORDER_COLOR)
add_textbox(slide, 5950000, 680000, 5600000, 330000,
            "关键特性", font_size=Pt(16), bold=True, color=SUBTITLE_BLUE)
features = [
    "GPU Direct Storage：GPU kernel 直接发起 NVMe 读请求",
    "Window Buffering：预取窗口缓冲，隐藏 SSD 延迟",
    "Storage Access Accumulator：批量合并 SSD 访问请求",
    "CPU Feature Buffer：热数据缓存（Zero-copy 映射）",
    "同构图 + 异构图双模式支持",
    "多 SSD 条带化（Page-level Striping）",
]
lines = [("• " + f, Pt(11), False, BODY_GRAY) for f in features]
add_multiline_textbox(slide, 5950000, 1050000, 5600000, 2100000, lines)

# 右侧下：性能对比
add_rect(slide, 5800000, 3400000, 5900000, 2800000, fill_color=rgb(0xFF,0xF5,0xE6), line_color=HIGHLIGHT_YEL)
add_textbox(slide, 5950000, 3480000, 5600000, 330000,
            "性能背景", font_size=Pt(16), bold=True, color=HIGHLIGHT_YEL)
perf_lines = [
    ("大规模 GNN 数据集规模：", Pt(12), True, TITLE_NAVY),
    ("IGB-Full: 547M 节点，1024维特征，~2TB", Pt(11), False, BODY_GRAY),
    ("OGB-Papers100M: 100M 节点，128维特征，~50GB", Pt(11), False, BODY_GRAY),
    ("MAG240M: 240M 节点，768维特征，~180GB", Pt(11), False, BODY_GRAY),
    ("", Pt(11), False, BODY_GRAY),
    ("GIDS 相比 mmap 基线：", Pt(12), True, TITLE_NAVY),
    ("数据加载吞吐量提升 2~5x（取决于 SSD 条数）", Pt(11), False, SUCCESS_GREEN),
    ("端到端训练时间缩短 40~60%", Pt(11), False, SUCCESS_GREEN),
]
add_multiline_textbox(slide, 5950000, 3850000, 5600000, 2200000, perf_lines)


# ── Slide 5: 章节分隔 - 系统架构 ──────────────────────────────
slide = prs.slides.add_slide(blank_layout)
add_section_divider(slide, "02", "系统架构总览",
                    "Four-Layer Architecture: Python → pybind11 → C++/CUDA → Hardware")


# ── Slide 6: 系统架构图 ────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
add_title_bar(slide, "系统架构总览 — 四层分解", "GIDS System Architecture Overview")

# 四层架构盒子
layers = [
    ("Layer 4  Python 应用层",
     "GIDS.py / GIDS_DGLDataLoader / 训练脚本",
     LIGHT_BG, SUBTITLE_BLUE, BORDER_COLOR),
    ("Layer 3  pybind11 绑定层",
     "BAM_Feature_Store.so (C++/CUDA → Python 桥接)",
     rgb(0xE8, 0xF8, 0xF5), SUCCESS_GREEN, rgb(0x27, 0xAE, 0x60)),
    ("Layer 2  CUDA/C++ 核心层",
     "gids_kernel.cu (GPU Kernels) + gids_nvme.cu (Host 管理逻辑)",
     rgb(0xEB, 0xF5, 0xFB), SUBTITLE_BLUE, SUBTITLE_BLUE),
    ("Layer 1  BaM 框架层",
     "libnvm.so: Controller / page_cache_t / bam_ptr<T> / NVMe Queue",
     rgb(0xF9, 0xEB, 0xEA), ERROR_RED, ERROR_RED),
    ("Layer 0  硬件层",
     "NVIDIA GPU (CUDA)  ←PCIe P2P→  NVMe SSD (/dev/libnvmX)",
     rgb(0x1A, 0x3A, 0x5C), ACCENT_CYAN, SECTION_BG),
]

box_h = 920000
start_y = 600000
for i, (title, desc, bg, title_c, border_c) in enumerate(layers):
    y = start_y + i * (box_h + 50000)
    add_rect(slide, 300000, y, SLIDE_W - 600000, box_h,
             fill_color=bg, line_color=border_c, line_width=Pt(2))
    if i == 4:
        add_textbox(slide, 450000, y + 100000, 9000000, 380000,
                    title, font_size=Pt(14), bold=True, color=ACCENT_CYAN)
        add_textbox(slide, 450000, y + 480000, 9000000, 350000,
                    desc, font_size=Pt(12), bold=False, color=rgb(0xBD, 0xD7, 0xEA))
    else:
        add_textbox(slide, 450000, y + 100000, 9000000, 380000,
                    title, font_size=Pt(14), bold=True, color=title_c)
        add_textbox(slide, 450000, y + 480000, 9000000, 350000,
                    desc, font_size=Pt(12), bold=False, color=BODY_GRAY)
    # 下箭头
    if i < 4:
        arrow_y = y + box_h
        add_textbox(slide, SLIDE_W//2 - 100000, arrow_y - 10000, 300000, 60000,
                    "▼", font_size=Pt(14), bold=True, color=ACCENT_CYAN, align=PP_ALIGN.CENTER)


# ── Slide 7: 章节分隔 - 核心组件 ──────────────────────────────
slide = prs.slides.add_slide(blank_layout)
add_section_divider(slide, "03", "核心组件深度解析",
                    "BaM Framework / GPU Kernels / BAM_Feature_Store / Python Interface")


# ── Slide 8: BaM 框架层 ────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
add_title_bar(slide, "核心组件 ① — BaM 框架层", "GPU-Initiated NVMe Access Framework")

# BaM 核心抽象表格
add_textbox(slide, 457200, 600000, SLIDE_W - 914400, 350000,
            "BaM (Block address Map) 核心抽象组件", font_size=Pt(16), bold=True, color=SUBTITLE_BLUE)

headers = ["组件", "类型", "核心作用", "在 GIDS 中的角色"]
rows = [
    ("Controller", "Host C++ 类", "管理一个 NVMe 设备，提供命令提交队列（SQ/CQ）", "每个 SSD 对应一个实例"),
    ("page_cache_t", "CUDA 设备/Host 类", "GPU 端 DRAM 页缓存，LRU 淘汰策略，管理页的分配/淘汰", "GPU 显存作为热数据二级缓存"),
    ("range_t<T>", "设备端模板类", "定义一段数据的逻辑地址→物理页映射（STRIPE/REPLICATE 分布）", "节点特征数据的地址空间"),
    ("array_t<T>", "设备端模板类", "可包含多个 range，支持多 SSD 条带化", "多盘特征数据的统一视图"),
    ("bam_ptr<T>", "GPU 智能指针", "operator[] 透明处理 page fault：miss→提交 NVMe 命令→等待 DMA 完成", "GPU kernel 内直接访问 SSD"),
    ("nvm_parallel_queue", "CUDA 设备类", "GPU kernel 内部 NVMe 命令队列，无需 CPU 介入", "异步 NVMe I/O 提交/完成"),
]

col_w = [1400000, 1600000, 4400000, 2500000]
col_x = [457200, 1857200, 3457200, 7857200]
row_h = 620000
header_y = 1050000
for ci, (hdr, cx, cw) in enumerate(zip(headers, col_x, col_w)):
    add_rect(slide, cx, header_y, cw - 30000, 380000, fill_color=SUBTITLE_BLUE)
    add_textbox(slide, cx + 60000, header_y + 80000, cw - 120000, 250000,
                hdr, font_size=Pt(12), bold=True, color=BG_WHITE)

for ri, row in enumerate(rows):
    y = header_y + 380000 + ri * row_h
    bg = LIGHT_BG if ri % 2 == 0 else BG_WHITE
    for ci, (cell, cx, cw) in enumerate(zip(row, col_x, col_w)):
        add_rect(slide, cx, y, cw - 30000, row_h - 20000,
                 fill_color=bg, line_color=BORDER_COLOR)
        col_c = SUBTITLE_BLUE if ci == 0 else BODY_GRAY
        bold = ci == 0
        add_textbox(slide, cx + 60000, y + 80000, cw - 120000, row_h - 160000,
                    cell, font_size=Pt(11), bold=bold, color=col_c)


# ── Slide 9: GPU Kernels ───────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
add_title_bar(slide, "核心组件 ② — GPU Kernel 层", "gids_kernel.cu 四种核心 Kernel")

kernel_data = [
    ("read_feature_kernel", "特征读取主核",
     ["输入：node IDs 数组 + bam_ptr<T> 引用",
      "输出：特征矩阵到 GPU output buffer",
      "实现：每线程读取一行特征向量",
      "关键：bam_ptr operator[] 触发 page fault 路径"],
     SUBTITLE_BLUE),
    ("cpu_buffer_kernel", "CPU 热数据分流核",
     ["检查 row_index < cpu_buffer_len",
      "热数据：从 CPU pinned buffer 读取（PCIe）",
      "冷数据：从 SSD 读取（bam_ptr GPU Direct）",
      "Zero-copy：cudaHostAllocMapped"],
     SUCCESS_GREEN),
    ("window_buffer_kernel", "预取窗口核",
     ["提前通知 BaM 页缓存即将访问的页面",
      "调用 bam_ptr.set_window_buffer_counter()",
      "增加目标页的预取优先级计数器",
      "在下一 batch 训练时数据已预热"],
     HIGHLIGHT_YEL),
    ("write_kernel", "特征写入核",
     ["将 GPU 内存中的 tensor 写回 SSD",
      "通过 bam_ptr operator[] 触发写操作",
      "用于 tensor_write.py 数据准备阶段",
      "支持多 SSD 条带化写入"],
     ERROR_RED),
]

box_w = (SLIDE_W - 1100000) // 4 - 20000
for i, (name, title, bullets, color) in enumerate(kernel_data):
    x = 457200 + i * (box_w + 35000)
    # 顶部颜色条
    add_rect(slide, x, 600000, box_w, 80000, fill_color=color)
    add_rect(slide, x, 680000, box_w, 5370000, fill_color=LIGHT_BG,
             line_color=BORDER_COLOR)
    add_textbox(slide, x + 60000, 730000, box_w - 120000, 350000,
                name, font_size=Pt(13), bold=True, color=color, wrap=False)
    add_textbox(slide, x + 60000, 1080000, box_w - 120000, 300000,
                title, font_size=Pt(12), bold=False, color=TITLE_NAVY, wrap=False)
    add_rect(slide, x + 60000, 1380000, box_w - 120000, 3, fill_color=BORDER_COLOR)
    lines = [("• " + b, Pt(11), False, BODY_GRAY) for b in bullets]
    add_multiline_textbox(slide, x + 60000, 1430000, box_w - 120000, 4000000, lines)

# 底部注意
add_rect(slide, 457200, 6150000, SLIDE_W - 914400, 350000,
         fill_color=rgb(0xFF, 0xF5, 0xE6), line_color=HIGHLIGHT_YEL)
add_textbox(slide, 600000, 6200000, SLIDE_W - 1200000, 280000,
            "⚠️  Corex 适配重点：<<<>>> kernel launch 语法需替换为 ixLaunchKernel()；bam_ptr 整体替换为 cuFile/IXFeatureStore 方案；cooperative_groups 依赖需补丁支持",
            font_size=Pt(11), bold=False, color=rgb(0x7D, 0x60, 0x08))


# ── Slide 10: BAM_Feature_Store ───────────────────────────────
slide = prs.slides.add_slide(blank_layout)
add_title_bar(slide, "核心组件 ③ — BAM_Feature_Store", "C++/CUDA 核心存储管理类（Host 端）")

# 左侧：类结构
add_rect(slide, 457200, 600000, 5200000, 5700000, fill_color=LIGHT_BG, line_color=BORDER_COLOR)
add_textbox(slide, 600000, 680000, 5000000, 320000,
            "BAM_Feature_Store<TYPE> 类成员", font_size=Pt(15), bold=True, color=SUBTITLE_BLUE)

struct_lines = [
    ("// BaM 核心对象", Pt(11), False, rgb(0x95, 0xA5, 0xA6)),
    ("page_cache_t *h_pc;       // Host 端页缓存句柄", Pt(11), False, SUBTITLE_BLUE),
    ("range_t<TYPE> *h_range;   // 数据范围映射", Pt(11), False, SUBTITLE_BLUE),
    ("array_t<TYPE> *a;         // GPU 端数组引用", Pt(11), False, SUBTITLE_BLUE),
    ("", Pt(11), False, BODY_GRAY),
    ("// CPU 缓冲优化", Pt(11), False, rgb(0x95, 0xA5, 0xA6)),
    ("GIDS_CPU_buffer<TYPE> CPU_buffer;", Pt(11), False, SUCCESS_GREEN),
    ("bool cpu_buffer_flag;     // 是否启用 CPU 缓冲", Pt(11), False, BODY_GRAY),
    ("bool seq_flag;            // 顺序/哈希模式", Pt(11), False, BODY_GRAY),
    ("", Pt(11), False, BODY_GRAY),
    ("// 控制参数", Pt(11), False, rgb(0x95, 0xA5, 0xA6)),
    ("uint32_t pageSize;        // 页大小 (默认 4096B)", Pt(11), False, BODY_GRAY),
    ("uint64_t numElems;        // 数据集总元素数", Pt(11), False, BODY_GRAY),
    ("uint32_t n_ctrls;         // SSD 数量", Pt(11), False, BODY_GRAY),
    ("", Pt(11), False, BODY_GRAY),
    ("// 核心方法", Pt(11), False, rgb(0x95, 0xA5, 0xA6)),
    ("void init_controllers(…); // 初始化 NVMe + 页缓存", Pt(11), False, HIGHLIGHT_YEL),
    ("void read_feature(…);     // 同构图特征读取", Pt(11), False, HIGHLIGHT_YEL),
    ("void read_feature_hetero(…); // 异构图多流并发", Pt(11), False, HIGHLIGHT_YEL),
    ("void read_feature_merged(…); // 批量合并读取", Pt(11), False, HIGHLIGHT_YEL),
    ("void cpu_backing_buffer(…);  // 分配 CPU Pinned", Pt(11), False, HIGHLIGHT_YEL),
    ("void set_cpu_buffer(…);      // 热节点预加载", Pt(11), False, HIGHLIGHT_YEL),
    ("void set_window_buffering(…);// 窗口预取", Pt(11), False, HIGHLIGHT_YEL),
    ("void store_tensor(…);         // 写回 SSD", Pt(11), False, HIGHLIGHT_YEL),
]
add_multiline_textbox(slide, 600000, 1050000, 5000000, 5000000, struct_lines,
                      font_name="Consolas")

# 右侧：pybind11 绑定
add_rect(slide, 5900000, 600000, 5800000, 2700000, fill_color=LIGHT_BG, line_color=BORDER_COLOR)
add_textbox(slide, 6050000, 680000, 5500000, 320000,
            "pybind11 Python 绑定", font_size=Pt(15), bold=True, color=SUBTITLE_BLUE)
bind_lines = [
    ("gids_nvme.cu 通过 pybind11 暴露：", Pt(12), True, TITLE_NAVY),
    ("", Pt(11), False, BODY_GRAY),
    ("• BAM_Feature_Store_float  — float32 特征", Pt(11), False, BODY_GRAY),
    ("• BAM_Feature_Store_half   — float16 特征", Pt(11), False, BODY_GRAY),
    ("• BAM_Feature_Store_long   — int64 节点 ID", Pt(11), False, BODY_GRAY),
    ("• GIDS_Controllers         — 控制器管理", Pt(11), False, BODY_GRAY),
    ("", Pt(11), False, BODY_GRAY),
    ("编译输出：BAM_Feature_Store.so", Pt(12), True, SUCCESS_GREEN),
    ("Python 侧 import：", Pt(12), True, TITLE_NAVY),
    ("  import BAM_Feature_Store as BFS", Pt(11), False, SUBTITLE_BLUE),
    ("  fs = BFS.BAM_Feature_Store_float()", Pt(11), False, SUBTITLE_BLUE),
]
add_multiline_textbox(slide, 6050000, 1050000, 5500000, 2100000, bind_lines,
                      font_name="Consolas")

# 右侧下：适配要点
add_rect(slide, 5900000, 3450000, 5800000, 2850000,
         fill_color=rgb(0xF9, 0xEB, 0xEA), line_color=ERROR_RED)
add_textbox(slide, 6050000, 3530000, 5500000, 320000,
            "Corex 适配要点", font_size=Pt(15), bold=True, color=ERROR_RED)
adapt_lines = [
    ("原版依赖（需替换）:", Pt(12), True, ERROR_RED),
    ("• bam_ptr<T> + page_cache_t + libnvm.so", Pt(11), False, BODY_GRAY),
    ("• /dev/libnvmX 裸 NVMe 字符设备", Pt(11), False, BODY_GRAY),
    ("• BaM 专有 page fault 机制", Pt(11), False, BODY_GRAY),
    ("", Pt(11), False, BODY_GRAY),
    ("适配方案（IXFeatureStore）:", Pt(12), True, SUCCESS_GREEN),
    ("• cuFileRead/Write 替代 bam_ptr.read()", Pt(11), False, BODY_GRAY),
    ("• 普通文件替代裸 NVMe 设备", Pt(11), False, BODY_GRAY),
    ("• OS page cache 替代 GPU 端 page_cache_t", Pt(11), False, BODY_GRAY),
    ("• 保留 CPU buffer / window buffer 策略", Pt(11), False, SUCCESS_GREEN),
]
add_multiline_textbox(slide, 6050000, 3900000, 5500000, 2200000, adapt_lines)


# ── Slide 11: 数据流分析 ───────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
add_title_bar(slide, "数据流分析 — 训练循环", "Training Loop Data Flow")

# 步骤流程
steps = [
    ("1", "DGL 采样", "MultiLayerNeighborSampler\n从大图 G 中采样子图 blocks\n输出：input_nodes (GPU tensor)"),
    ("2", "GIDS 特征读取", "GIDS.fetch_feature(dim, it, device)\n→ BAM_FS.read_feature()\n→ GPU Kernel 访问 bam_ptr"),
    ("3", "BaM 页缓存", "page_cache_t 查询\nCache Hit：直接从 GPU 显存返回\nCache Miss：提交 NVMe 命令"),
    ("4", "NVMe DMA", "GPU 直接发起 NVMe 命令\n通过 PCIe P2P DMA\n数据从 SSD 写入 GPU 显存"),
    ("5", "GNN 训练", "SAGEConv / GCN / GAT\n使用特征 tensor 前向传播\n反向传播更新参数"),
]

arrow_x = [457200, 2700000, 5000000, 7300000, 9600000]
box_w = 2100000
for i, (num, title, desc) in enumerate(steps):
    x = arrow_x[i]
    add_rect(slide, x, 700000, box_w, 4000000, fill_color=LIGHT_BG, line_color=BORDER_COLOR)
    add_rect(slide, x, 700000, box_w, 280000, fill_color=ACCENT_CYAN)
    add_textbox(slide, x + 80000, 720000, box_w - 160000, 240000,
                f"Step {num}", font_size=Pt(12), bold=True, color=BG_WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, x + 80000, 1040000, box_w - 160000, 380000,
                title, font_size=Pt(13), bold=True, color=SUBTITLE_BLUE, align=PP_ALIGN.CENTER)
    add_textbox(slide, x + 80000, 1500000, box_w - 160000, 2900000,
                desc, font_size=Pt(11), bold=False, color=BODY_GRAY, align=PP_ALIGN.CENTER)

    if i < 4:
        add_textbox(slide, x + box_w + 20000, 1600000, 200000, 300000,
                    "→", font_size=Pt(20), bold=True, color=ACCENT_CYAN, align=PP_ALIGN.CENTER)

# 流说明
add_rect(slide, 457200, 5000000, SLIDE_W - 914400, 1500000,
         fill_color=rgb(0xE8, 0xF8, 0xF5), line_color=SUCCESS_GREEN)
add_textbox(slide, 600000, 5080000, SLIDE_W - 1200000, 280000,
            "CUDA Stream 并行化", font_size=Pt(14), bold=True, color=SUCCESS_GREEN)
stream_lines = [
    "Stream[0]: [GNN Forward+Backward]  [GNN Forward+Backward]  [GNN Forward+Backward]  ...",
    "Stream[1]: [Window Prefetch]        [Window Prefetch]        [Window Prefetch]        ...",
    "Stream[2]: [Feature Read (cuFile)]  [Feature Read (cuFile)]  [Feature Read (cuFile)]  ...",
]
for i, line in enumerate(stream_lines):
    color = [TITLE_NAVY, HIGHLIGHT_YEL, SUBTITLE_BLUE][i]
    add_textbox(slide, 600000, 5380000 + i * 330000, SLIDE_W - 1200000, 300000,
                line, font_size=Pt(11), bold=False, color=color, font_name="Consolas")


# ── Slide 12: 章节分隔 - 三大优化 ─────────────────────────────
slide = prs.slides.add_slide(blank_layout)
add_section_divider(slide, "04", "三大优化策略",
                    "Window Buffering / Storage Access Accumulator / CPU Feature Buffer")


# ── Slide 13: 三大优化策略 ────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
add_title_bar(slide, "三大优化策略", "Three Core Optimization Strategies")

strategies = [
    ("01  Window Buffering", SUBTITLE_BLUE,
     "问题：GPU kernel 访问 SSD 时，page fault 导致高延迟",
     [
         "在处理当前 batch N 时，提前通知 BaM 缓存 batch N+1 的页面",
         "GPU kernel：set_window_buffer_counter() 增加目标页预取计数",
         "Stream[1] 异步执行，不阻塞 Stream[0] 训练流",
         "效果：将 SSD 延迟隐藏在计算时间内，接近零延迟",
     ],
     "延迟隐藏 / 流水线重叠"),
    ("02  Storage Access Accumulator", SUCCESS_GREEN,
     "问题：小批量下每次 SSD 访问都有固定延迟开销（效率低）",
     [
         "根据 SSD 带宽/延迟计算最优合并数量 N",
         "累积 N 个 batch 的节点 ID 后，一次性合并提交",
         "调用 read_feature_merged()，单次 DMA 传输完成",
         "效果：NVMe 命令数减少 N 倍，提升带宽利用率",
     ],
     "批量合并 / 减少 I/O 次数"),
    ("03  CPU Feature Buffer", HIGHLIGHT_YEL,
     "问题：高频访问节点（热数据）每次都从 SSD 读取，浪费带宽",
     [
         "PageRank 算法识别 Top-K 高度数热节点",
         "cudaHostAllocMapped 分配 CPU Pinned Memory",
         "cudaHostGetDevicePointer 获取 GPU 侧虚拟地址",
         "kernel 内检查 row_index < cpu_len 分流：热→CPU，冷→SSD",
     ],
     "Zero-copy 热数据 / 减少 SSD 访问"),
]

for i, (title, color, problem, bullets, effect) in enumerate(strategies):
    y = 680000 + i * 1830000
    add_rect(slide, 457200, y, SLIDE_W - 914400, 1700000,
             fill_color=LIGHT_BG, line_color=color, line_width=Pt(2))
    add_rect(slide, 457200, y, 350000, 1700000, fill_color=color)
    add_textbox(slide, 900000, y + 80000, 4000000, 360000,
                title, font_size=Pt(15), bold=True, color=color)
    add_textbox(slide, 900000, y + 440000, 4000000, 320000,
                problem, font_size=Pt(11), bold=False, color=BODY_GRAY)
    lines = [("• " + b, Pt(11), False, BODY_GRAY) for b in bullets]
    add_multiline_textbox(slide, 5200000, y + 80000, 6500000, 1550000, lines)
    add_rect(slide, SLIDE_W - 1400000, y + 200000, 1100000, 350000, fill_color=color)
    add_textbox(slide, SLIDE_W - 1380000, y + 250000, 1060000, 270000,
                effect, font_size=Pt(10), bold=True, color=BG_WHITE, align=PP_ALIGN.CENTER)


# ── Slide 14: 章节分隔 - Corex 适配挑战 ──────────────────────
slide = prs.slides.add_slide(blank_layout)
add_section_divider(slide, "05", "Corex 适配挑战分析",
                    "Hardware Differences / Software Dependencies / Migration Risk")


# ── Slide 15: 适配挑战 ────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
add_title_bar(slide, "Corex 适配核心挑战", "Key Adaptation Challenges")

# 对比表格
add_textbox(slide, 457200, 600000, SLIDE_W - 914400, 320000,
            "NVIDIA 原版 vs Iluvatar Corex 关键差异", font_size=Pt(15), bold=True, color=SUBTITLE_BLUE)

dims = ["GPU 运行时", "编译器", "GDS 用户态库", "GDS 内核驱动", "存储设备访问",
        "Kernel Launch 语法", "PyTorch 设备名", "DGL 支持"]
nvidia_vals = [
    "CUDA Runtime / Driver",
    "nvcc (NVIDIA CUDA Compiler)",
    "libcufile.so (NVIDIA GDS)",
    "nvidia-fs.ko / /dev/nvidia-fs",
    "BaM: /dev/libnvmX (裸 NVMe)",
    "<<<grid, block>>> 语法",
    "cuda:0",
    "pip install dgl (CUDA 版)"
]
corex_vals = [
    "IX Runtime + CUDA 兼容映射层",
    "ixc (Iluvatar CUDA Compiler)",
    "Corex libcufile.so (兼容 API)",
    "itrfs.ko / /dev/itrfs ✅ 已加载",
    "IXFeatureStore + cuFile + 普通文件",
    "ixLaunchKernel() API 调用",
    "cuda:0 ✅ (不使用 ix:0)",
    "需源码编译 Corex CUDA 版 DGL",
]
status = ["ok", "warn", "ok", "ok", "ok", "warn", "ok", "warn"]
status_text = ["已兼容", "需修改", "已兼容", "已验证", "方案确定", "需适配", "已兼容", "进行中"]

col_w2 = [2200000, 3500000, 3500000, 1500000]
col_x2 = [457200, 2657200, 6157200, 9657200]
headers2 = ["维度", "NVIDIA 原版", "Corex 适配版", "状态"]

h_y = 1030000
for ci, (hdr, cx, cw) in enumerate(zip(headers2, col_x2, col_w2)):
    add_rect(slide, cx, h_y, cw - 20000, 350000, fill_color=SUBTITLE_BLUE)
    add_textbox(slide, cx + 60000, h_y + 70000, cw - 120000, 230000,
                hdr, font_size=Pt(12), bold=True, color=BG_WHITE)

row_h2 = 600000
for ri in range(len(dims)):
    y = h_y + 350000 + ri * row_h2
    bg = LIGHT_BG if ri % 2 == 0 else BG_WHITE
    st = status[ri]
    badge_colors = {"ok": SUCCESS_GREEN, "warn": HIGHLIGHT_YEL, "err": ERROR_RED}
    row_data = [dims[ri], nvidia_vals[ri], corex_vals[ri], status_text[ri]]
    for ci, (cell, cx, cw) in enumerate(zip(row_data, col_x2, col_w2)):
        add_rect(slide, cx, y, cw - 20000, row_h2 - 15000,
                 fill_color=bg, line_color=BORDER_COLOR)
        if ci == 3:
            add_rect(slide, cx + 60000, y + 100000, cw - 180000, 350000,
                     fill_color=badge_colors[st])
            add_textbox(slide, cx + 80000, y + 150000, cw - 200000, 240000,
                        cell, font_size=Pt(11), bold=True, color=BG_WHITE, align=PP_ALIGN.CENTER)
        else:
            col_c = SUBTITLE_BLUE if ci == 0 else BODY_GRAY
            add_textbox(slide, cx + 60000, y + 80000, cw - 120000, row_h2 - 200000,
                        cell, font_size=Pt(11), bold=(ci == 0), color=col_c)


# ── Slide 16: 章节分隔 - 适配方案 ────────────────────────────
slide = prs.slides.add_slide(blank_layout)
add_section_divider(slide, "06", "适配方案与当前进展",
                    "cuFile Replaces BaM / IXFeatureStore Design / Current Status")


# ── Slide 17: 三方案对比 ──────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
add_title_bar(slide, "适配方案选型对比", "Solution Architecture Decision")

add_textbox(slide, 457200, 600000, SLIDE_W - 914400, 320000,
            "三种方案权衡分析", font_size=Pt(15), bold=True, color=SUBTITLE_BLUE)

plans = [
    ("方案 A ★ 当前采用",
     "cuFile 替代 BaM",
     ACCENT_CYAN,
     ["用 IXFeatureStore + cuFileRead/Write 读写普通文件",
      "依赖 Corex 官方 GDS 栈（itrfs.ko + libcufile.so）",
      "Python 接口保持不变（GIDS_IX.py 替代 GIDS.py）",
      "存储从裸设备改为普通文件路径"],
     ["工作量小，稳定可调试", "依赖 Corex 官方 GDS", "OS page cache 透明",
      "快速验证端到端路径"],
     ["无 GPU kernel 内透明 page fault", "bam_ptr 语义改变", "理论延迟略高于 BaM"],
     "✅ 当前采用"),
    ("方案 B",
     "BaM 完整移植",
     rgb(0x95, 0xA5, 0xA6),
     ["移植 BaM 裸 NVMe、GPU page cache、NVMe SQ/CQ",
      "实现 /dev/libnvmX 等价的 Corex 字符设备",
      "移植 page_cache_t GPU 端页缓存",
      "保留 bam_ptr<T> 透明 page fault 机制"],
     ["理论延迟最低", "保留原论文路径", "GPU kernel 透明访问"],
     ["工作量大（>3 个月）", "硬件依赖多（MMIO/P2P）",
      "调试验证困难", "风险高"],
     "⏸ 暂不采用"),
    ("方案 C",
     "POSIX 降级模式",
     rgb(0xBD, 0xC3, 0xC7),
     ["cuFileRead 失败时自动降级",
      "使用 pread/pwrite + CPU pinned buffer",
      "ixMemcpy H2D 传输到 GPU",
      "作为 cuFile GDS 路径的 fallback"],
     ["兼容性最强", "无需内核驱动", "调试方便"],
     ["吞吐量低（2~3x 差距）", "CPU 再次成为瓶颈", "不适合生产环境"],
     "🔄 作为 Fallback"),
]

for i, (name, subtitle, color, impl, pros, cons, verdict) in enumerate(plans):
    x = 457200 + i * 3900000
    w = 3750000
    border = color if i == 0 else rgb(0xCC, 0xCC, 0xCC)
    bg = LIGHT_BG if i == 0 else BG_WHITE
    lw = Pt(3) if i == 0 else Pt(1.5)
    add_rect(slide, x, 1000000, w, 5300000, fill_color=bg, line_color=border, line_width=lw)
    add_rect(slide, x, 1000000, w, 380000, fill_color=color)
    add_textbox(slide, x + 80000, 1040000, w - 160000, 300000,
                name, font_size=Pt(14), bold=True, color=BG_WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, x + 80000, 1430000, w - 160000, 300000,
                subtitle, font_size=Pt(13), bold=True, color=color, align=PP_ALIGN.CENTER)
    add_textbox(slide, x + 80000, 1780000, w - 160000, 260000,
                "实现方式：", font_size=Pt(12), bold=True, color=TITLE_NAVY)
    impl_lines = [("  • " + p, Pt(11), False, BODY_GRAY) for p in impl]
    add_multiline_textbox(slide, x + 80000, 2040000, w - 160000, 1200000, impl_lines)
    add_textbox(slide, x + 80000, 3290000, w - 160000, 260000,
                "优点：", font_size=Pt(12), bold=True, color=SUCCESS_GREEN)
    pros_lines = [("  ✓ " + p, Pt(11), False, SUCCESS_GREEN) for p in pros]
    add_multiline_textbox(slide, x + 80000, 3550000, w - 160000, 800000, pros_lines)
    add_textbox(slide, x + 80000, 4400000, w - 160000, 260000,
                "缺点：", font_size=Pt(12), bold=True, color=ERROR_RED)
    cons_lines = [("  ✗ " + c, Pt(11), False, ERROR_RED) for c in cons]
    add_multiline_textbox(slide, x + 80000, 4660000, w - 160000, 800000, cons_lines)
    add_rect(slide, x + 80000, 5550000, w - 160000, 580000, fill_color=color)
    add_textbox(slide, x + 80000, 5630000, w - 160000, 420000,
                verdict, font_size=Pt(14), bold=True, color=BG_WHITE, align=PP_ALIGN.CENTER)


# ── Slide 18: 适配后架构图 ────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
add_title_bar(slide, "适配后架构 — IXFeatureStore", "Adapted Architecture: GIDS-IX on Corex Platform")

# 左侧：原架构（变灰）
add_rect(slide, 300000, 650000, 4200000, 5800000, fill_color=BG_WHITE, line_color=rgb(0xCC,0xCC,0xCC))
add_textbox(slide, 400000, 700000, 4000000, 320000,
            "原架构 (NVIDIA + BaM)", font_size=Pt(14), bold=True, color=rgb(0x95,0xA5,0xA6))
orig_layers = [
    ("GIDS.py + GIDS_DGLDataLoader", rgb(0xCC,0xCC,0xCC)),
    ("BAM_Feature_Store.so (pybind11)", rgb(0xCC,0xCC,0xCC)),
    ("gids_kernel.cu (GPU Kernels)", rgb(0xCC,0xCC,0xCC)),
    ("BaM: bam_ptr + page_cache_t", ERROR_RED),
    ("/dev/libnvmX (裸 NVMe 设备)", ERROR_RED),
    ("NVIDIA GPU + NVMe SSD (P2P)", rgb(0xCC,0xCC,0xCC)),
]
for i, (text, color) in enumerate(orig_layers):
    y = 1100000 + i * 780000
    add_rect(slide, 400000, y, 3900000, 680000, fill_color=rgb(0xF5,0xF5,0xF5), line_color=rgb(0xCC,0xCC,0xCC))
    add_textbox(slide, 500000, y + 100000, 3700000, 500000,
                text, font_size=Pt(12), bold=False, color=color, align=PP_ALIGN.CENTER)

# 中间箭头
add_textbox(slide, 4600000, 3200000, 600000, 500000,
            "⟹", font_size=Pt(36), bold=True, color=ACCENT_CYAN, align=PP_ALIGN.CENTER)
add_textbox(slide, 4550000, 3700000, 700000, 300000,
            "Corex\n适配", font_size=Pt(12), bold=True, color=ACCENT_CYAN, align=PP_ALIGN.CENTER)

# 右侧：新架构
add_rect(slide, 5350000, 650000, 6400000, 5800000, fill_color=LIGHT_BG, line_color=ACCENT_CYAN, line_width=Pt(2))
add_textbox(slide, 5450000, 700000, 6200000, 320000,
            "适配架构 (Corex + cuFile)", font_size=Pt(14), bold=True, color=ACCENT_CYAN)
new_layers = [
    ("GIDS_IX.py + GIDS_DGLDataLoader", SUCCESS_GREEN, "Python 接口兼容保留"),
    ("IXFeatureStore.so (pybind11)", SUCCESS_GREEN, "替代 BAM_Feature_Store"),
    ("ix_feature_store.cu (Host 端)", SUBTITLE_BLUE, "无需 GPU Direct Page Fault"),
    ("cuFileRead / cuFileWrite", ACCENT_CYAN, "标准 GDS API"),
    ("/mnt/nvme0/node_feat.bin", SUCCESS_GREEN, "普通文件替代裸设备"),
    ("Corex GPU + itrfs.ko + /dev/itrfs", SUCCESS_GREEN, "已验证可用"),
]
for i, (text, color, note) in enumerate(new_layers):
    y = 1100000 + i * 780000
    add_rect(slide, 5450000, y, 6100000, 680000, fill_color=BG_WHITE, line_color=color, line_width=Pt(1.5))
    add_textbox(slide, 5550000, y + 80000, 5900000, 340000,
                text, font_size=Pt(12), bold=True, color=color, align=PP_ALIGN.CENTER)
    add_textbox(slide, 5550000, y + 420000, 5900000, 230000,
                note, font_size=Pt(10), bold=False, color=BODY_GRAY, align=PP_ALIGN.CENTER)


# ── Slide 19: 章节分隔 - 关键适配工作 ────────────────────────
slide = prs.slides.add_slide(blank_layout)
add_section_divider(slide, "07", "关键适配工作清单",
                    "DGL Compilation / cooperative_groups Patch / cuFile Integration")


# ── Slide 20: 关键适配工作 1 ──────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
add_title_bar(slide, "关键适配工作 — DGL CUDA 版编译", "DGL v1.1.3 Corex CUDA Build")

add_textbox(slide, 457200, 600000, SLIDE_W - 914400, 320000,
            "为什么必须编译 Corex CUDA 版 DGL？", font_size=Pt(15), bold=True, color=SUBTITLE_BLUE)

# 原因
add_rect(slide, 457200, 1000000, 5800000, 2200000, fill_color=rgb(0xF9,0xEB,0xEA), line_color=ERROR_RED)
add_textbox(slide, 600000, 1080000, 5600000, 320000,
            "pip install dgl 的问题", font_size=Pt(14), bold=True, color=ERROR_RED)
reason_lines = [
    ("• pip 安装的是 CPU 版 libdgl.so", Pt(12), False, BODY_GRAY),
    ("• 不包含 CUDA Device API 实现", Pt(12), False, BODY_GRAY),
    ("• pin_memory_() 和 UVA 采样不可用 → GIDS DataLoader 崩溃", Pt(12), True, ERROR_RED),
    ("• graph.to('cuda:0') 调用失败", Pt(12), True, ERROR_RED),
]
add_multiline_textbox(slide, 600000, 1430000, 5600000, 1500000, reason_lines)

# 6个编译修复
add_textbox(slide, 457200, 3350000, SLIDE_W - 914400, 320000,
            "DGL v1.1.3 Corex 编译 6 项关键修复", font_size=Pt(15), bold=True, color=SUBTITLE_BLUE)

fixes = [
    ("-Xcompiler 格式", "nvcc 接受逗号分隔，ixc 不接受\n→ 改为 separate_arguments 空格列表", "warn"),
    ("fp16.cuh 重定义", "DGL 与 Corex SDK 的 fp16.cuh 冲突\n→ 使用 #pragma once 保护", "warn"),
    ("CCCL 禁用", "Corex 不依赖 NVIDIA CCCL 头文件\n→ -DUSE_CCCL=OFF 编译参数", "warn"),
    ("omp.h 路径", "OpenMP 头文件路径不同\n→ 修正 include 搜索路径", "ok"),
    ("array_iterator.h", "DGL 依赖 CUDA Thrust 迭代器\n→ 使用 Corex 兼容版", "warn"),
    ("gpu_cache 禁用", "cooperative_groups tiled_partition 触发\nLLVM %laneid PTX → 暂时禁用 gpu_cache", "err"),
]

fbox_w = (SLIDE_W - 1200000) // 3
for i, (title, desc, status) in enumerate(fixes):
    col = i % 3
    row = i // 3
    x = 457200 + col * (fbox_w + 60000)
    y = 3750000 + row * 1200000
    status_colors = {"ok": SUCCESS_GREEN, "warn": HIGHLIGHT_YEL, "err": ERROR_RED}
    c = status_colors[status]
    add_rect(slide, x, y, fbox_w, 1100000, fill_color=LIGHT_BG, line_color=c, line_width=Pt(2))
    add_rect(slide, x, y, fbox_w, 80000, fill_color=c)
    add_textbox(slide, x + 80000, y + 100000, fbox_w - 160000, 300000,
                title, font_size=Pt(13), bold=True, color=c)
    add_textbox(slide, x + 80000, y + 430000, fbox_w - 160000, 600000,
                desc, font_size=Pt(11), bold=False, color=BODY_GRAY)


# ── Slide 21: 关键适配工作 2 ──────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
add_title_bar(slide, "关键适配工作 — cooperative_groups 补丁", "SWPM-918-gids 补丁分析")

# 问题描述
add_rect(slide, 457200, 600000, 5800000, 2300000, fill_color=rgb(0xF9,0xEB,0xEA), line_color=ERROR_RED)
add_textbox(slide, 600000, 680000, 5600000, 320000,
            "问题：meta_group_rank() / %laneid", font_size=Pt(15), bold=True, color=ERROR_RED)
prob_lines = [
    ("DGL gpu_cache 和 GIDS window_buffer 使用了:", Pt(12), True, TITLE_NAVY),
    ("  cooperative_groups::tiled_partition<32>()", Pt(12), False, SUBTITLE_BLUE),
    ("  meta_group_rank() / meta_group_size()", Pt(12), False, SUBTITLE_BLUE),
    ("", Pt(11), False, BODY_GRAY),
    ("问题根因：", Pt(12), True, ERROR_RED),
    ("  ixc LLVM 后端在翻译 tiled_partition 时生成 %laneid PTX 寄存器", Pt(11), False, BODY_GRAY),
    ("  %laneid 在 Corex GPU 架构中未定义/不支持", Pt(11), False, BODY_GRAY),
    ("  非源码级问题，是编译器后端（LLVM llc）层面", Pt(11), False, BODY_GRAY),
]
add_multiline_textbox(slide, 600000, 1030000, 5600000, 1700000, prob_lines)

# 解决方案
add_rect(slide, 457200, 3000000, 5800000, 2300000, fill_color=rgb(0xE8,0xF8,0xF5), line_color=SUCCESS_GREEN)
add_textbox(slide, 600000, 3080000, 5600000, 320000,
            "SWPM-918-gids 分支补丁方案", font_size=Pt(15), bold=True, color=SUCCESS_GREEN)
sol_lines = [
    ("补丁位置：", Pt(12), True, TITLE_NAVY),
    ("  /home/corex/sw_home_1/sw_home/sdk/ixdriver/include/cooperative_groups.h", Pt(11), False, SUBTITLE_BLUE),
    ("", Pt(11), False, BODY_GRAY),
    ("补丁内容：", Pt(12), True, TITLE_NAVY),
    ("  meta_group_rank() → 返回固定值（API 层绕过）", Pt(11), False, BODY_GRAY),
    ("  meta_group_size() → 返回固定值", Pt(11), False, BODY_GRAY),
    ("", Pt(11), False, BODY_GRAY),
    ("重要限制：", Pt(12), True, HIGHLIGHT_YEL),
    ("  仅对 API 调用层有效，不解决编译器后端 %laneid 问题", Pt(11), False, BODY_GRAY),
    ("  gpu_cache 短期内仍需禁用（不影响 GIDS 核心功能）", Pt(11), False, BODY_GRAY),
]
add_multiline_textbox(slide, 600000, 3430000, 5600000, 1700000, sol_lines)

# 右侧：影响范围
add_rect(slide, 6500000, 600000, 5200000, 4700000, fill_color=LIGHT_BG, line_color=BORDER_COLOR)
add_textbox(slide, 6650000, 680000, 5000000, 320000,
            "影响范围分析", font_size=Pt(15), bold=True, color=SUBTITLE_BLUE)
impact_lines = [
    ("受影响组件：", Pt(13), True, ERROR_RED),
    ("  ❌ DGL gpu_cache 模块", Pt(12), False, ERROR_RED),
    ("  ❌ GIDS window_buffer kernel", Pt(12), False, HIGHLIGHT_YEL),
    ("  ❌ 任何使用 tiled_partition 的 kernel", Pt(12), False, BODY_GRAY),
    ("", Pt(11), False, BODY_GRAY),
    ("不受影响组件：", Pt(13), True, SUCCESS_GREEN),
    ("  ✅ IXFeatureStore 核心读写逻辑", Pt(12), False, SUCCESS_GREEN),
    ("  ✅ cuFileRead/Write API", Pt(12), False, SUCCESS_GREEN),
    ("  ✅ GIDS Python 接口层", Pt(12), False, SUCCESS_GREEN),
    ("  ✅ DGL 图采样（CPU 侧）", Pt(12), False, SUCCESS_GREEN),
    ("", Pt(11), False, BODY_GRAY),
    ("短期应对：", Pt(13), True, SUBTITLE_BLUE),
    ("  禁用 gpu_cache，不影响 GIDS 核心", Pt(12), False, BODY_GRAY),
    ("  window_buffer 在 Python 层完成调度", Pt(12), False, BODY_GRAY),
    ("", Pt(11), False, BODY_GRAY),
    ("长期方向：", Pt(13), True, SUBTITLE_BLUE),
    ("  推动 SWPM-918 正式合入 ixc LLVM 后端", Pt(12), False, BODY_GRAY),
    ("  或提供不依赖 %laneid 的替代实现", Pt(12), False, BODY_GRAY),
]
add_multiline_textbox(slide, 6650000, 1050000, 5000000, 4000000, impact_lines)


# ── Slide 22: 关键适配工作 3 ──────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
add_title_bar(slide, "关键适配工作 — cuFile/GDS 集成", "cuFile API Integration & GDS Stack")

# cuFile API 表格
add_textbox(slide, 457200, 600000, SLIDE_W - 914400, 320000,
            "标准 cuFile API 对应关系", font_size=Pt(15), bold=True, color=SUBTITLE_BLUE)

cufile_apis = [
    ("cuFileDriverOpen()", "初始化 cuFile 驱动", "依赖 itrfs.ko + /dev/itrfs", "✅ 已验证"),
    ("cuFileHandleRegister()", "注册文件句柄（关联 fd）", "打开后立即注册", "✅ 已验证"),
    ("cuFileBufRegister()", "注册 GPU buffer", "GPU malloc 后注册，加速 DMA 路径", "✅ 已验证"),
    ("cuFileRead(fh, buf, size, file_off, buf_off)", "GPU DMA 读取", "直接从 NVMe→GPU，绕过 CPU", "✅ 已验证"),
    ("cuFileWrite(fh, buf, size, file_off, buf_off)", "GPU DMA 写入", "用于 tensor_write.py 数据准备", "✅ 已验证"),
    ("cuFileBufDeregister()", "注销 GPU buffer", "释放前必须注销", "✅ 已验证"),
    ("cuFileHandleDeregister()", "注销文件句柄", "关闭前注销", "✅ 已验证"),
    ("cuFileDriverClose()", "关闭 cuFile 驱动", "程序退出前调用", "✅ 已验证"),
]

headers3 = ["API", "功能", "说明", "状态"]
col_w3 = [2800000, 1600000, 4000000, 1300000]
col_x3 = [457200, 3257200, 4857200, 8857200]

h_y3 = 1030000
for ci, (hdr, cx, cw) in enumerate(zip(headers3, col_x3, col_w3)):
    add_rect(slide, cx, h_y3, cw - 20000, 350000, fill_color=SUBTITLE_BLUE)
    add_textbox(slide, cx + 60000, h_y3 + 70000, cw - 120000, 230000,
                hdr, font_size=Pt(12), bold=True, color=BG_WHITE)

for ri, row in enumerate(cufile_apis):
    y = h_y3 + 350000 + ri * 560000
    bg = LIGHT_BG if ri % 2 == 0 else BG_WHITE
    for ci, (cell, cx, cw) in enumerate(zip(row, col_x3, col_w3)):
        add_rect(slide, cx, y, cw - 20000, 530000, fill_color=bg, line_color=BORDER_COLOR)
        c = SUBTITLE_BLUE if ci == 0 else (SUCCESS_GREEN if ci == 3 else BODY_GRAY)
        bold = ci in (0, 3)
        fn = "Consolas" if ci == 0 else "微软雅黑"
        add_textbox(slide, cx + 60000, y + 80000, cw - 120000, 400000,
                    cell, font_size=Pt(11), bold=bold, color=c, font_name=fn)

# 降级机制
add_rect(slide, 457200, 5680000, SLIDE_W - 914400, 900000,
         fill_color=rgb(0xFF,0xF5,0xE6), line_color=HIGHLIGHT_YEL)
add_textbox(slide, 600000, 5760000, SLIDE_W - 1200000, 280000,
            "自动降级机制（itrfs.ko 不可用时）", font_size=Pt(13), bold=True, color=HIGHLIGHT_YEL)
add_textbox(slide, 600000, 6050000, SLIDE_W - 1200000, 440000,
            "cuFileDriverOpen() 失败 → 检测到 /dev/itrfs 不可用 → 自动切换 POSIX pread()/pwrite() 路径 → ixMemcpy H2D 传输 → 保证程序不崩溃（性能降低约 2~3x）",
            font_size=Pt(11), bold=False, color=rgb(0x7D,0x60,0x08))


# ── Slide 23: 依赖环境清单 ────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
add_title_bar(slide, "依赖环境清单", "Environment Dependencies & Versions")

# 左侧：Corex SDK 包
add_rect(slide, 457200, 600000, 5400000, 5800000, fill_color=LIGHT_BG, line_color=BORDER_COLOR)
add_textbox(slide, 600000, 680000, 5200000, 320000,
            "Corex 4.5.0 SDK 组件", font_size=Pt(15), bold=True, color=SUBTITLE_BLUE)

sdk_items = [
    ("ixdriver", "libcudart.so + libcufile.so + libcuda.so", "✅"),
    ("ixinfer", "libcuinfer.so.7", "✅"),
    ("ixpti", "libcupti.so (CUPTI Profiling)", "✅"),
    ("ixblas", "libcublas.so (BLAS)", "✅"),
    ("ixdnn", "libcudnn.so (DNN)", "✅"),
    ("ixsparse", "libcusparse.so (Sparse)", "✅"),
    ("ixfft", "libcufft.so (FFT)", "✅"),
    ("ixccl", "libnccl.so (Collective Comm)", "✅"),
    ("itrfs", "itrfs.ko + /dev/itrfs (GDS 驱动)", "✅"),
]

for i, (pkg, desc, status) in enumerate(sdk_items):
    y = 1080000 + i * 520000
    bg = LIGHT_BG if i % 2 == 0 else BG_WHITE
    add_rect(slide, 600000, y, 5100000, 480000, fill_color=bg, line_color=BORDER_COLOR)
    add_textbox(slide, 650000, y + 80000, 1200000, 330000,
                pkg, font_size=Pt(12), bold=True, color=SUBTITLE_BLUE, font_name="Consolas")
    add_textbox(slide, 1900000, y + 80000, 3000000, 330000,
                desc, font_size=Pt(11), bold=False, color=BODY_GRAY)
    add_textbox(slide, 4900000, y + 80000, 700000, 330000,
                status, font_size=Pt(13), bold=True, color=SUCCESS_GREEN, align=PP_ALIGN.CENTER)

# 右侧：Python 包
add_rect(slide, 6100000, 600000, 5600000, 5800000, fill_color=LIGHT_BG, line_color=BORDER_COLOR)
add_textbox(slide, 6250000, 680000, 5400000, 320000,
            "Python 依赖包", font_size=Pt(15), bold=True, color=SUBTITLE_BLUE)

py_items = [
    ("torch", "2.10.0+corex.4.5.0", "Iluvatar 定制 PyTorch", "✅"),
    ("dgl", "1.1.3 CUDA 版", "需源码编译 Corex 版", "🔧"),
    ("pybind11", "2.13.6", "C++/Python 绑定编译", "✅"),
    ("numpy", ">=1.14.0", "数值计算", "✅"),
    ("ogb", "1.3.6", "Open Graph Benchmark 数据集", "✅"),
    ("torchdata", "0.7.1", "注意：0.11.0 不兼容，必须降级", "⚠️"),
    ("cooperative_groups.h", "SWPM-918-gids", "补丁分支版本", "✅"),
]

for i, (pkg, ver, desc, status) in enumerate(py_items):
    y = 1080000 + i * 660000
    bg = LIGHT_BG if i % 2 == 0 else BG_WHITE
    c_status = SUCCESS_GREEN if status == "✅" else (HIGHLIGHT_YEL if status == "⚠️" else SUBTITLE_BLUE)
    add_rect(slide, 6250000, y, 5300000, 620000, fill_color=bg, line_color=BORDER_COLOR)
    add_textbox(slide, 6300000, y + 60000, 1400000, 250000,
                pkg, font_size=Pt(12), bold=True, color=SUBTITLE_BLUE, font_name="Consolas")
    add_textbox(slide, 6300000, y + 330000, 1400000, 220000,
                ver, font_size=Pt(11), bold=False, color=BODY_GRAY, font_name="Consolas")
    add_textbox(slide, 7750000, y + 60000, 3400000, 450000,
                desc, font_size=Pt(11), bold=False, color=BODY_GRAY)
    add_textbox(slide, 11000000, y + 150000, 400000, 300000,
                status, font_size=Pt(16), bold=True, color=c_status, align=PP_ALIGN.CENTER)


# ── Slide 24: 章节分隔 - 后期工作流程 ────────────────────────
slide = prs.slides.add_slide(blank_layout)
add_section_divider(slide, "08", "后期工作流程与里程碑",
                    "End-to-End Validation Path / Phase Goals / Timeline")


# ── Slide 25: 当前进展 ────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
add_title_bar(slide, "当前进展状态", "Current Status (2026-06-12)")

milestones = [
    ("✅ 架构分析完成", SUCCESS_GREEN,
     ["GIDS 全部源码深度分析（gids_kernel.cu / gids_nvme.cu）",
      "BaM 框架与 Corex 兼容性逐 API 对比完成",
      "三种适配方案（cuFile/BaM移植/POSIX）权衡决策完成",
      "选定方案 A（cuFile 替代 BaM）"]),
    ("✅ 核心代码移植完成", SUCCESS_GREEN,
     ["ix_feature_store.h/cu (~450行) 新增实现",
      "GIDS_IX.py 适配层完成",
      "build_ix.sh 编译脚本完成",
      "run.sh 一键自动化脚本完成"]),
    ("✅ 编译验证通过", SUCCESS_GREEN,
     ["sandbox 环境编译通过（g++ 模式）",
      "Docker 真机 Corex 环境编译通过",
      "pybind11 模块 IXFeatureStore 可正常 import",
      "DGL v1.1.3 主库 Corex 编译进行中（6项修复已应用）"]),
    ("✅ 基础设施就绪", SUCCESS_GREEN,
     ["itrfs.ko 已加载，/dev/itrfs 已就绪",
      "cooperative_groups.h SWPM-918-gids 补丁已部署",
      "Corex 4.5.0 SDK 全部组件安装验证",
      "GDS 加速路径 cuFile API 基础调用验证通过"]),
    ("🔧 DGL CUDA 版编译", HIGHLIGHT_YEL,
     ["6项编译问题已修复并应用",
      "主库编译进度：90%（gpu_cache 模块禁用绕过）",
      "预计完成：DGL 源码编译 + Python 包安装",
      "验证：dgl.rand_graph().to('cuda:0') 正常"]),
    ("⏳ 端到端训练验证", SUBTITLE_BLUE,
     ["前置：DGL CUDA 版安装完成",
      "准备数据：tensor_write_ix.py 写入特征到 NVMe",
      "验证 GIDS_unit_test.py 单元测试",
      "全链路：homogenous_train_ix.py 运行"]),
]

box_w3 = (SLIDE_W - 1100000) // 3
for i, (title, color, items) in enumerate(milestones):
    col = i % 3
    row = i // 3
    x = 457200 + col * (box_w3 + 50000)
    y = 680000 + row * 2800000
    add_rect(slide, x, y, box_w3, 2600000, fill_color=LIGHT_BG, line_color=color, line_width=Pt(2))
    add_rect(slide, x, y, box_w3, 340000, fill_color=color)
    add_textbox(slide, x + 80000, y + 60000, box_w3 - 160000, 250000,
                title, font_size=Pt(13), bold=True, color=BG_WHITE)
    lines = [("• " + item, Pt(11), False, BODY_GRAY) for item in items]
    add_multiline_textbox(slide, x + 80000, y + 400000, box_w3 - 160000, 2100000, lines)


# ── Slide 26: 后期工作流程 ────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
add_title_bar(slide, "后期工作流程", "Future Work Pipeline & Milestones")

phases = [
    ("Phase 1\n近期 1~2 周",
     ACCENT_CYAN,
     "DGL CUDA 版完成 + 端到端验证",
     [
         "完成 DGL v1.1.3 Corex CUDA 源码编译",
         "验证 graph.pin_memory_() 和 UVA 采样",
         "准备测试数据：tensor_write_ix.py 写入特征",
         "执行 GIDS_unit_test.py 单元测试",
         "完成 homogenous_train_ix.py 端到端训练",
         "测量基础吞吐量：Epoch时间 / DataLoader时间",
     ]),
    ("Phase 2\n1~2 个月",
     SUBTITLE_BLUE,
     "性能优化 + 多 SSD 支持",
     [
         "多 SSD 条带化：应用层 page_idx % n_ctrls 分片",
         "Window Buffering 策略在新架构下的优化",
         "Storage Access Accumulator 调参优化",
         "CPU Feature Buffer（PageRank 热节点）集成",
         "cuFile batch API 探索（减少 per-request 开销）",
         "Profiling：ixpti + nsys 定位瓶颈",
     ]),
    ("Phase 3\n2~3 个月",
     HIGHLIGHT_YEL,
     "异构图 + 高级特性",
     [
         "heterogeneous_train_ix.py 异构图支持",
         "多流并发读取（Stream[0]/[1]/[2] 流水线）",
         "cooperative_groups tiled_partition 正式修复",
         "gpu_cache 模块重新启用（SWPM-918 合入后）",
         "ClusterGCN 采样器适配",
         "大规模数据集：IGB-Full / MAG240M 验证",
     ]),
    ("Phase 4\n持续优化",
     SUCCESS_GREEN,
     "生产化 + 文档完善",
     [
         "GNN 训练基准测试（对比 NVIDIA 原版 GIDS）",
         "Corex 平台性能调优指南",
         "多 GPU / 多机分布式 GDS 扩展",
         "cuFile 高级 API（cuFileStreamRead）探索",
         "与 ixdriver 深度集成（专用 GDS 优化路径）",
         "开源社区回馈与文档完善",
     ]),
]

ph_w = (SLIDE_W - 1100000) // 4
for i, (phase, color, goal, tasks) in enumerate(phases):
    x = 457200 + i * (ph_w + 50000)
    # 顶部宽条
    add_rect(slide, x, 680000, ph_w, 500000, fill_color=color)
    add_textbox(slide, x + 60000, 700000, ph_w - 120000, 460000,
                phase, font_size=Pt(14), bold=True, color=BG_WHITE, align=PP_ALIGN.CENTER)
    # 目标
    add_rect(slide, x, 1180000, ph_w, 380000, fill_color=LIGHT_BG, line_color=color)
    add_textbox(slide, x + 60000, 1210000, ph_w - 120000, 340000,
                goal, font_size=Pt(11), bold=True, color=color, align=PP_ALIGN.CENTER)
    # 任务列表
    add_rect(slide, x, 1560000, ph_w, 4800000, fill_color=LIGHT_BG, line_color=BORDER_COLOR)
    lines = [(f"{'▸' if j == 0 else '•'} {t}", Pt(12) if j == 0 else Pt(11),
              j == 0, color if j == 0 else BODY_GRAY)
             for j, t in enumerate(tasks)]
    add_multiline_textbox(slide, x + 80000, 1620000, ph_w - 160000, 4650000, lines)


# ── Slide 27: 风险与对策 ──────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
add_title_bar(slide, "风险评估与对策", "Risk Assessment & Mitigation")

risks = [
    ("itrfs.ko 稳定性",
     "高",
     "Corex GDS 驱动稳定性低于 NVIDIA nvidia-fs.ko，可能在高并发 I/O 下出现问题",
     "建立 POSIX fallback 自动切换机制；定期 dmesg 监控驱动日志；联系 Iluvatar 内核团队"),
    ("cooperative_groups 后端问题",
     "中",
     "%laneid PTX 在 ixc LLVM 后端无法翻译，影响 gpu_cache 和 window_buffer kernel",
     "短期禁用 gpu_cache；window_buffer 降级到 Python 层；跟进 SWPM-918 工单"),
    ("DGL CUDA 版兼容性",
     "中",
     "DGL 依赖大量 CUDA 特有 API，Corex 兼容层可能有遗漏",
     "逐模块测试，禁用不兼容子模块；优先保证采样核心功能"),
    ("cuFile GDS 吞吐量",
     "中",
     "Corex cuFile 实现可能不如 NVIDIA 原版，影响 SSD 读取带宽",
     "Profiling 对比；探索 cuFile batch read；多 SSD 条带化补偿"),
    ("PyTorch 版本耦合",
     "低",
     "torch-2.10.0+corex.4.5.0 与 SDK 4.5.0 强耦合，版本升级风险",
     "锁定版本；升级前完整回归测试；维护版本兼容矩阵"),
    ("大规模数据集准备",
     "低",
     "IGB-Full (~2TB) 写入 NVMe 需要专用数据准备流程",
     "tensor_write_ix.py 批量写入；验证多 SSD 写入正确性"),
]

row_h4 = 900000
for i, (risk, level, desc, mitigation) in enumerate(risks):
    col = i % 2
    row = i // 2
    x = 457200 + col * 5750000
    y = 680000 + row * row_h4
    level_color = {"高": ERROR_RED, "中": HIGHLIGHT_YEL, "低": SUCCESS_GREEN}[level]
    add_rect(slide, x, y, 5550000, row_h4 - 30000, fill_color=LIGHT_BG, line_color=BORDER_COLOR)
    add_rect(slide, x, y, 5550000, 80000, fill_color=level_color)
    add_textbox(slide, x + 80000, y + 110000, 3500000, 300000,
                risk, font_size=Pt(13), bold=True, color=TITLE_NAVY)
    add_rect(slide, x + 3900000, y + 100000, 1400000, 260000, fill_color=level_color)
    add_textbox(slide, x + 3950000, y + 140000, 1300000, 200000,
                f"风险：{level}", font_size=Pt(11), bold=True, color=BG_WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, x + 80000, y + 430000, 5300000, 250000,
                "▶ " + desc, font_size=Pt(11), bold=False, color=BODY_GRAY)
    add_textbox(slide, x + 80000, y + 650000, 5300000, 250000,
                "✓ " + mitigation, font_size=Pt(11), bold=False, color=SUCCESS_GREEN)


# ── Slide 28: 总结 ─────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill_color=SECTION_BG)
add_rect(slide, 0, 0, 220000, SLIDE_H, fill_color=ACCENT_CYAN)
add_rect(slide, 0, SLIDE_H - 100000, SLIDE_W, 100000, fill_color=ACCENT_CYAN)

add_textbox(slide, 500000, 400000, SLIDE_W - 1000000, 600000,
            "总结 Summary", font_size=Pt(36), bold=True, color=ACCENT_CYAN)

summary_points = [
    ("架构理解", ACCENT_CYAN,
     "GIDS 是四层架构的 GPU Direct Storage 系统，核心是 BaM 框架提供 GPU 端透明页故障 NVMe 访问能力"),
    ("适配策略", rgb(0x58, 0xD6, 0xA8),
     "选择方案 A（cuFile 替代 BaM）：用 IXFeatureStore + Corex libcufile.so 替代 BaM 裸 NVMe，降低移植风险"),
    ("当前状态", rgb(0xFF, 0xC0, 0x60),
     "核心代码移植完成（ix_feature_store.cu），编译验证通过，GDS 基础设施就绪，DGL 编译进行中"),
    ("后期路径", rgb(0xFF, 0x80, 0x80),
     "Phase 1 端到端验证 → Phase 2 性能优化 → Phase 3 异构图 → Phase 4 生产化"),
]

for i, (title, color, desc) in enumerate(summary_points):
    y = 1100000 + i * 1200000
    add_rect(slide, 500000, y, SLIDE_W - 1000000, 1100000,
             fill_color=rgb(0x1F, 0x45, 0x70), line_color=color, line_width=Pt(2))
    add_rect(slide, 500000, y, 350000, 1100000, fill_color=color)
    add_textbox(slide, 1000000, y + 100000, 3000000, 400000,
                f"0{i+1}  {title}", font_size=Pt(16), bold=True, color=color)
    add_textbox(slide, 1000000, y + 540000, SLIDE_W - 1600000, 500000,
                desc, font_size=Pt(13), bold=False, color=rgb(0xBD, 0xD7, 0xEA))

add_textbox(slide, 500000, 6000000, SLIDE_W - 1000000, 500000,
            "Iluvatar GPU Platform  |  GIDS Adaptation  |  Corex 4.5.0  |  2026-06",
            font_size=Pt(13), bold=False, color=ACCENT_CYAN, align=PP_ALIGN.CENTER)


# ─── 保存 ──────────────────────────────────────────────────────
output_path = "/root/GIDS_cufile/GIDS-架构分析与Corex适配-PPT.pptx"
prs.save(output_path)
print(f"✅ PPT 已生成：{output_path}")
print(f"   共 {len(prs.slides)} 张幻灯片")
